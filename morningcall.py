import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import math
import os
import re
from PIL import Image
import tempfile


# Configurações visuais do PowerPoint
FONTE_TITULO = "Raleway"
FONTE_CORPO = "Raleway"
COR_TITULO = RGBColor(0, 40, 60)
COR_CORPO = RGBColor(40, 80, 90)

# Valores de referência calibrados para A4 vertical (8.264 x 11.694 in). O
# template Morningcall.pptx pode vir em outro tamanho de slide (ex.: o layout
# ampliado usado no analise.pptx) — por isso todos escalam proporcionalmente
# ao tamanho real do template via _aplicar_escala_pagina(), chamada uma vez
# ao carregar a Presentation. Sem essa chamada, os valores abaixo são usados
# como estão (equivalente a um slide A4).
_REF_LARGURA = Inches(8.264)
_REF_ALTURA = Inches(11.694)

_REF_TAMANHO_TITULO = 22
_REF_TAMANHO_CORPO = 11
_REF_ESPACAMENTO_LINHA_CORPO = 17
_REF_ESPACO_APOS_PARAGRAFO = 12
_REF_TAMANHO_LOGO_FALLBACK = 10

_REF_MARGEM_X = Inches(0.83)
_REF_TOPO_CONTEUDO = Inches(0.90)
_REF_LARGURA_CONTEUDO = Inches(6.60)
_REF_ALTURA_TITULO = Inches(0.40)
_REF_ESPACO_TITULO_CORPO = Inches(0.10)
_REF_ESPACO_ENTRE_PARAGRAFOS = Inches(0.14)
_REF_ESPACO_ENTRE_SECOES = Inches(0.32)
_REF_RODAPE_LIMITE = Inches(10.55)

_REF_ALTURA_LINHA_CORPO = Inches(0.225)
_REF_ALTURA_MINIMA_CORPO = Inches(0.2)

_REF_LOGO_MAX_W = Inches(0.55)
_REF_LOGO_MAX_H = Inches(0.55)
_REF_LOGO_MARGEM_DIREITA = Inches(0.80)
_REF_LOGO_Y = Inches(10.55)

_REF_ALTURA_MAX_GRAFICO = Inches(9.0)
_REF_ESPACO_TEXTO_GRAFICO = Inches(0.25)
_REF_ESPACO_APOS_GRAFICO = Inches(0.35)

CARACTERES_POR_LINHA = 88
LINHAS_EXTRAS_POR_PARAGRAFO = 0.75

LOGO_RODAPE = "icone_final.png"

TAMANHO_TITULO = Pt(_REF_TAMANHO_TITULO)
TAMANHO_CORPO = Pt(_REF_TAMANHO_CORPO)
ESPACAMENTO_LINHA_CORPO = Pt(_REF_ESPACAMENTO_LINHA_CORPO)
ESPACO_APOS_PARAGRAFO = Pt(_REF_ESPACO_APOS_PARAGRAFO)
TAMANHO_LOGO_FALLBACK = Pt(_REF_TAMANHO_LOGO_FALLBACK)

MARGEM_X = _REF_MARGEM_X
TOPO_CONTEUDO = _REF_TOPO_CONTEUDO
LARGURA_CONTEUDO = _REF_LARGURA_CONTEUDO
ALTURA_TITULO = _REF_ALTURA_TITULO
ESPACO_TITULO_CORPO = _REF_ESPACO_TITULO_CORPO
ESPACO_ENTRE_PARAGRAFOS = _REF_ESPACO_ENTRE_PARAGRAFOS
ESPACO_ENTRE_SECOES = _REF_ESPACO_ENTRE_SECOES
RODAPE_LIMITE = _REF_RODAPE_LIMITE

ALTURA_LINHA_CORPO = _REF_ALTURA_LINHA_CORPO
ALTURA_MINIMA_CORPO = _REF_ALTURA_MINIMA_CORPO

LOGO_MAX_W = _REF_LOGO_MAX_W
LOGO_MAX_H = _REF_LOGO_MAX_H
LOGO_MARGEM_DIREITA = _REF_LOGO_MARGEM_DIREITA
LOGO_Y = _REF_LOGO_Y

ALTURA_MAX_GRAFICO = _REF_ALTURA_MAX_GRAFICO
ESPACO_TEXTO_GRAFICO = _REF_ESPACO_TEXTO_GRAFICO
ESPACO_APOS_GRAFICO = _REF_ESPACO_APOS_GRAFICO


def _aplicar_escala_pagina(prs):
    """Reescala as constantes de layout para o tamanho real do slide do template.

    Sempre recalcula a partir dos valores _REF_* (referência A4), nunca a
    partir das constantes públicas já escaladas — assim a função pode ser
    chamada mais de uma vez (ex.: várias gerações na mesma sessão do
    Streamlit) sem acumular escala.
    """
    global MARGEM_X, TOPO_CONTEUDO, LARGURA_CONTEUDO, ALTURA_TITULO
    global ESPACO_TITULO_CORPO, ESPACO_ENTRE_PARAGRAFOS, ESPACO_ENTRE_SECOES, RODAPE_LIMITE
    global ALTURA_LINHA_CORPO, ALTURA_MINIMA_CORPO, LOGO_MAX_W, LOGO_MAX_H, LOGO_MARGEM_DIREITA, LOGO_Y
    global ALTURA_MAX_GRAFICO, ESPACO_TEXTO_GRAFICO, ESPACO_APOS_GRAFICO
    global TAMANHO_TITULO, TAMANHO_CORPO, ESPACAMENTO_LINHA_CORPO, ESPACO_APOS_PARAGRAFO
    global TAMANHO_LOGO_FALLBACK

    escala_w = prs.slide_width / _REF_LARGURA
    escala_h = prs.slide_height / _REF_ALTURA

    MARGEM_X = int(_REF_MARGEM_X * escala_w)
    LARGURA_CONTEUDO = int(_REF_LARGURA_CONTEUDO * escala_w)
    LOGO_MARGEM_DIREITA = int(_REF_LOGO_MARGEM_DIREITA * escala_w)
    LOGO_MAX_W = int(_REF_LOGO_MAX_W * escala_w)

    TOPO_CONTEUDO = int(_REF_TOPO_CONTEUDO * escala_h)
    ALTURA_TITULO = int(_REF_ALTURA_TITULO * escala_h)
    ESPACO_TITULO_CORPO = int(_REF_ESPACO_TITULO_CORPO * escala_h)
    ESPACO_ENTRE_PARAGRAFOS = int(_REF_ESPACO_ENTRE_PARAGRAFOS * escala_h)
    ESPACO_ENTRE_SECOES = int(_REF_ESPACO_ENTRE_SECOES * escala_h)
    RODAPE_LIMITE = int(_REF_RODAPE_LIMITE * escala_h)
    ALTURA_LINHA_CORPO = int(_REF_ALTURA_LINHA_CORPO * escala_h)
    ALTURA_MINIMA_CORPO = int(_REF_ALTURA_MINIMA_CORPO * escala_h)
    LOGO_MAX_H = int(_REF_LOGO_MAX_H * escala_h)
    LOGO_Y = int(_REF_LOGO_Y * escala_h)
    ALTURA_MAX_GRAFICO = int(_REF_ALTURA_MAX_GRAFICO * escala_h)
    ESPACO_TEXTO_GRAFICO = int(_REF_ESPACO_TEXTO_GRAFICO * escala_h)
    ESPACO_APOS_GRAFICO = int(_REF_ESPACO_APOS_GRAFICO * escala_h)

    TAMANHO_TITULO = Pt(max(6, round(_REF_TAMANHO_TITULO * escala_h)))
    TAMANHO_CORPO = Pt(max(6, round(_REF_TAMANHO_CORPO * escala_h)))
    ESPACAMENTO_LINHA_CORPO = Pt(max(6, round(_REF_ESPACAMENTO_LINHA_CORPO * escala_h)))
    ESPACO_APOS_PARAGRAFO = Pt(max(6, round(_REF_ESPACO_APOS_PARAGRAFO * escala_h)))
    TAMANHO_LOGO_FALLBACK = Pt(max(6, round(_REF_TAMANHO_LOGO_FALLBACK * escala_h)))
def carregar_graficos_disponiveis(data_inicio_semana=None, data_fim_semana=None):
    """
    Carrega os gráficos disponíveis a partir do arquivo graficos.py.

    Formato recomendado em graficos.py:

    def gerar_graficos():
        return {
            "Gráfico Ibovespa": "graficos/ibovespa.png",
            "Gráfico S&P 500": "graficos/sp500.png",
        }

    O valor pode ser:
    - caminho de imagem: .png, .jpg, .jpeg;
    - figura matplotlib;
    - figura plotly.
    """
    try:
        import graficos
    except Exception:
        return {}

    if hasattr(graficos, "gerar_graficos"):
        try:
            return graficos.gerar_graficos(
                data_inicio_semana=data_inicio_semana,
                data_fim_semana=data_fim_semana,
            )
        except Exception as e:
            st.warning(f"Não foi possível carregar os gráficos de graficos.py: {e}")
            return {}

    if hasattr(graficos, "obter_graficos"):
        try:
            return graficos.obter_graficos()
        except Exception as e:
            st.warning(f"Não foi possível carregar os gráficos de graficos.py: {e}")
            return {}

    return {}


def converter_grafico_para_imagem(grafico):
    """
    Converte diferentes tipos de gráfico em um caminho de imagem para inserir no PowerPoint.
    """
    if grafico is None:
        return None

    if isinstance(grafico, str):
        return grafico if os.path.exists(grafico) else None

    # Matplotlib figure
    if hasattr(grafico, "savefig"):
        arquivo_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        grafico.savefig(arquivo_temp.name, dpi=200, bbox_inches="tight")
        return arquivo_temp.name

    # Plotly figure
    if hasattr(grafico, "write_image"):
        arquivo_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        grafico.write_image(arquivo_temp.name)
        return arquivo_temp.name

    return None


def estimar_altura_grafico(caminho_grafico):
    if not caminho_grafico or not os.path.exists(caminho_grafico):
        return 0

    with Image.open(caminho_grafico) as img:
        largura_px, altura_px = img.size

    proporcao = altura_px / largura_px
    altura = int(LARGURA_CONTEUDO * proporcao)

    return min(altura, ALTURA_MAX_GRAFICO)


def adicionar_grafico(slide, caminho_grafico, y):
    """Adiciona gráfico abaixo do texto, mantendo proporção e respeitando a largura do conteúdo."""
    if not caminho_grafico or not os.path.exists(caminho_grafico):
        return y

    with Image.open(caminho_grafico) as img:
        largura_px, altura_px = img.size

    proporcao = altura_px / largura_px
    altura = int(LARGURA_CONTEUDO * proporcao)
    altura = min(altura, ALTURA_MAX_GRAFICO)

    slide.shapes.add_picture(
        caminho_grafico,
        MARGEM_X,
        y,
        width=LARGURA_CONTEUDO,
        height=altura
    )

    return y + altura + ESPACO_APOS_GRAFICO


def adicionar_titulo(slide, texto, y):
    if y < TOPO_CONTEUDO:
        y = TOPO_CONTEUDO
    caixa = slide.shapes.add_textbox(MARGEM_X, y, LARGURA_CONTEUDO, ALTURA_TITULO)
    caixa.name = "Caixa de texto - título"
    tf = caixa.text_frame
    tf.clear()
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run = p.add_run()
    run.text = texto
    run.font.name = FONTE_TITULO
    run.font.size = TAMANHO_TITULO
    run.font.bold = True
    run.font.italic = False
    run.font.color.rgb = COR_TITULO

    return y + ALTURA_TITULO


def adicionar_corpo(slide, texto, y, altura):
    altura = min(altura, max(ALTURA_MINIMA_CORPO, RODAPE_LIMITE - y))
    caixa = slide.shapes.add_textbox(MARGEM_X, y, LARGURA_CONTEUDO, altura)
    caixa.name = "Caixa de texto - corpo"
    tf = caixa.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    paragrafos = [p.strip() for p in texto.split("\n") if p.strip()] if texto else [""]

    for i, paragrafo in enumerate(paragrafos):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.JUSTIFY
        p.line_spacing = ESPACAMENTO_LINHA_CORPO
        p.space_after = ESPACO_APOS_PARAGRAFO

        for j, seg in enumerate(re.split(r'\*([^*]+)\*', paragrafo)):
            if not seg:
                continue
            run = p.add_run()
            run.text = seg
            run.font.name = FONTE_CORPO
            run.font.size = TAMANHO_CORPO
            run.font.color.rgb = COR_CORPO
            run.font.bold = (j % 2 == 1)

    return y + altura


def limpar_slide(slide):
    """Remove todos os elementos soltos do slide, deixando a página pronta para receber conteúdo."""
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)


def adicionar_logo_rodape(slide):
    """Adiciona a logo da Your Capital no canto inferior direito, mantendo a proporção original."""
    if os.path.exists(LOGO_RODAPE):
        with Image.open(LOGO_RODAPE) as img:
            largura_px, altura_px = img.size

        proporcao = largura_px / altura_px

        largura = LOGO_MAX_W
        altura = int(largura / proporcao)

        if altura > LOGO_MAX_H:
            altura = LOGO_MAX_H
            largura = int(altura * proporcao)

        slide_width = slide.part.slide_layout.part.package.presentation_part.presentation.slide_width
        logo_x = slide_width - LOGO_MARGEM_DIREITA - largura

        slide.shapes.add_picture(
            LOGO_RODAPE,
            logo_x,
            LOGO_Y,
            width=largura,
            height=altura
        )
    else:
        # Fallback: caso o arquivo da logo não esteja na pasta, escreve o nome no rodapé.
        caixa = slide.shapes.add_textbox(Inches(5.65), Inches(10.70), Inches(1.8), Inches(0.25))
        tf = caixa.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = "YOUR CAPITAL"
        run.font.name = FONTE_TITULO
        run.font.size = TAMANHO_LOGO_FALLBACK
        run.font.bold = True
        run.font.color.rgb = COR_TITULO


# --- Helpers para preenchimento de slides dinâmicos ---

def estimar_linhas(texto):
    if not texto:
        return 1

    total_linhas = 0
    paragrafos = [p.strip() for p in texto.split("\n") if p.strip()]

    for paragrafo in paragrafos:
        tamanho = len(paragrafo)
        linhas_paragrafo = max(1, math.ceil(tamanho / CARACTERES_POR_LINHA))
        total_linhas += linhas_paragrafo

    # Ajuste do espaço entre parágrafos.
    # O PowerPoint não quebra automaticamente para outro slide; por isso
    # a estimativa precisa ser conservadora para evitar overflow da caixa.
    total_linhas += max(0, len(paragrafos) * LINHAS_EXTRAS_POR_PARAGRAFO)

    return total_linhas


def estimar_altura_corpo(texto):
    return ALTURA_LINHA_CORPO * estimar_linhas(texto)


def dividir_texto_por_altura(texto, altura_disponivel):
    """
    Divide o texto preservando parágrafos. A primeira parte cabe na altura disponível;
    a segunda continua no próximo slide.
    """
    max_linhas = max(1, int(altura_disponivel / ALTURA_LINHA_CORPO))
    paragrafos = [p.strip() for p in texto.split("\n") if p.strip()]

    parte_1 = []
    parte_2 = []
    linhas_usadas = 0
    passou_limite = False

    for paragrafo in paragrafos:
        palavras = paragrafo.split()
        linhas_paragrafo = []
        linha_atual = ""

        for palavra in palavras:
            tentativa = palavra if not linha_atual else f"{linha_atual} {palavra}"
            if len(tentativa) <= CARACTERES_POR_LINHA:
                linha_atual = tentativa
            else:
                linhas_paragrafo.append(linha_atual)
                linha_atual = palavra

        if linha_atual:
            linhas_paragrafo.append(linha_atual)

        linhas_necessarias = max(1, len(linhas_paragrafo)) + LINHAS_EXTRAS_POR_PARAGRAFO
        if parte_1:
            linhas_necessarias += 0.15

        if not passou_limite and linhas_usadas + linhas_necessarias <= max_linhas:
            parte_1.append(paragrafo)
            linhas_usadas += linhas_necessarias
        else:
            passou_limite = True
            parte_2.append(paragrafo)

    if not parte_1 and parte_2:
        # Se um único parágrafo for maior que a página, divide por palavras.
        paragrafo = parte_2.pop(0)
        palavras = paragrafo.split()
        linhas_usadas = 0
        trecho_1 = []
        trecho_2 = []
        linha_atual = ""

        for palavra in palavras:
            tentativa = palavra if not linha_atual else f"{linha_atual} {palavra}"
            if len(tentativa) <= CARACTERES_POR_LINHA:
                linha_atual = tentativa
            else:
                linhas_usadas += 1
                if linhas_usadas <= max_linhas:
                    trecho_1.extend(linha_atual.split())
                else:
                    trecho_2.extend(linha_atual.split())
                linha_atual = palavra

        if linha_atual:
            linhas_usadas += 1
            if linhas_usadas <= max_linhas:
                trecho_1.extend(linha_atual.split())
            else:
                trecho_2.extend(linha_atual.split())

        parte_1 = [" ".join(trecho_1)]
        resto = " ".join(trecho_2)
        parte_2 = ([resto] if resto else []) + parte_2

    return "\n".join(parte_1), "\n".join(parte_2)


def criar_slide_conteudo(prs):
    """
    Cria uma página de conteúdo do zero: slide em branco + logo no rodapé.
    Assim o layout fica totalmente controlado pelo código.
    """
    layout_em_branco = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout_em_branco)
    adicionar_logo_rodape(slide)
    return slide


def mover_slide_para_antes_do_fechamento(prs, slide):
    """
    Move o slide recém-criado para antes da última página, preservando a página final de fechamento.
    """
    sld_id_lst = prs.slides._sldIdLst
    slide_id = slide.slide_id
    ids = list(sld_id_lst)

    elemento_slide = None
    for elemento in ids:
        if int(elemento.get("id")) == slide_id:
            elemento_slide = elemento
            break

    if elemento_slide is not None:
        sld_id_lst.remove(elemento_slide)
        sld_id_lst.insert(len(sld_id_lst) - 1, elemento_slide)


def obter_novo_slide_conteudo(prs):
    novo_slide = criar_slide_conteudo(prs)
    mover_slide_para_antes_do_fechamento(prs, novo_slide)
    return novo_slide


def adicionar_secao_no_relatorio(prs, slide_atual, titulo, texto, y_atual, grafico=None):
    # A seção deve começar na página atual sempre que houver espaço mínimo
    # para o título e pelo menos algumas linhas do corpo. O restante continua
    # automaticamente na próxima página, se necessário.
    altura_minima_inicio = ALTURA_TITULO + ESPACO_TITULO_CORPO + (ALTURA_LINHA_CORPO * 3)

    if y_atual > TOPO_CONTEUDO and y_atual + altura_minima_inicio > RODAPE_LIMITE:
        slide_atual = obter_novo_slide_conteudo(prs)
        y_atual = TOPO_CONTEUDO

    if titulo.strip():
        y_atual = adicionar_titulo(slide_atual, titulo, y_atual)
        if texto.strip():
            y_atual += Inches(0.14)

    while texto:
        altura_disponivel = RODAPE_LIMITE - y_atual

        if altura_disponivel <= ALTURA_LINHA_CORPO:
            slide_atual = obter_novo_slide_conteudo(prs)
            y_atual = TOPO_CONTEUDO
            altura_disponivel = RODAPE_LIMITE - y_atual

        altura_corpo = estimar_altura_corpo(texto)

        if altura_corpo <= altura_disponivel:
            # Usa a altura estimada, sem expandir demais a caixa visual.
            # A folga contra overflow é tratada pela estimativa conservadora acima.
            altura_caixa = min(altura_corpo, altura_disponivel)
            y_atual = adicionar_corpo(slide_atual, texto, y_atual, altura_caixa)
            texto = ""
        else:
            parte_1, parte_2 = dividir_texto_por_altura(texto, altura_disponivel)
            if parte_1.strip():
                y_atual = adicionar_corpo(slide_atual, parte_1, y_atual, min(altura_disponivel, estimar_altura_corpo(parte_1)))
            texto = parte_2

            if texto:
                slide_atual = obter_novo_slide_conteudo(prs)
                y_atual = TOPO_CONTEUDO

    caminho_grafico = converter_grafico_para_imagem(grafico)

    if caminho_grafico:
        altura_grafico = estimar_altura_grafico(caminho_grafico)
        altura_necessaria = ESPACO_TEXTO_GRAFICO + altura_grafico + ESPACO_APOS_GRAFICO

        if y_atual + altura_necessaria > RODAPE_LIMITE:
            slide_atual = obter_novo_slide_conteudo(prs)
            y_atual = TOPO_CONTEUDO
        else:
            y_atual += ESPACO_TEXTO_GRAFICO

        y_atual = adicionar_grafico(slide_atual, caminho_grafico, y_atual)

    y_atual += ESPACO_ENTRE_SECOES
    return slide_atual, y_atual


def show_morningcall():
    import datetime

    _logo = Image.open("logo_final.png")
    st.image(_logo, use_container_width=False, width=800)

    st.title("Gerador de Relatório")

    st.markdown(
        "Preencha o título e o texto de cada seção. O PowerPoint será gerado a partir do template, mantendo a capa e preenchendo as páginas em branco."
    )

    # --- Semana de referência para o Painel de Mercado ---
    with st.expander("Semana de referência (Painel de Mercado)", expanded=True):
        _hoje = datetime.date.today()
        _dia = _hoje.weekday()                            # 0=seg, 6=dom
        _dias_desde_sexta = (_dia - 4) % 7 or 7          # dias corridos até a última sexta
        _default_fim = _hoje - datetime.timedelta(days=_dias_desde_sexta)
        _default_ini = _default_fim - datetime.timedelta(days=4)

        col_ini, col_fim = st.columns(2)
        with col_ini:
            data_inicio = st.date_input(
                "Início da semana",
                value=_default_ini,
                key="semana_inicio",
            )
        with col_fim:
            data_fim = st.date_input(
                "Fim da semana",
                value=_default_fim,
                key="semana_fim",
            )

        if data_inicio > data_fim:
            st.error("A data de início deve ser anterior à data de fim.")
            st.stop()

    graficos_disponiveis = carregar_graficos_disponiveis(
        data_inicio_semana=data_inicio,
        data_fim_semana=data_fim,
    )
    opcoes_graficos = ["Sem gráfico"] + list(graficos_disponiveis.keys())

    quantidade_secoes = st.number_input(
        "Quantidade de seções",
        min_value=1,
        max_value=20,
        value=4,
        step=1
    )

    titulos_padrao = [
        "Análise Macroeconômica",
        "Análise Carteira de Investimentos",
        "Fundos de Renda Fixa Crédito Privado",
        "Movimentações",
    ]

    secoes = []

    for i in range(int(quantidade_secoes)):
        st.markdown(f"### Seção {i + 1}")

        titulo_padrao = titulos_padrao[i] if i < len(titulos_padrao) else f"Seção {i + 1}"

        titulo = st.text_input(
            f"Título da seção {i + 1}",
            titulo_padrao,
            key=f"titulo_secao_{i}"
        )

        texto = st.text_area(
            f"Texto da seção {i + 1}",
            height=300,
            key=f"texto_secao_{i}"
        )

        grafico_escolhido = st.selectbox(
            f"Gráfico da seção {i + 1}",
            opcoes_graficos,
            index=0,
            key=f"grafico_secao_{i}"
        )

        grafico = None if grafico_escolhido == "Sem gráfico" else graficos_disponiveis.get(grafico_escolhido)

        secoes.append((titulo, texto, grafico))

    if st.button("Gerar PowerPoint"):
        prs = Presentation("Morningcall.pptx")

        # Reescala margens, fontes e espaçamentos para o tamanho real do slide
        # do template (pode não ser mais A4 — ver _aplicar_escala_pagina).
        _aplicar_escala_pagina(prs)

        # Estrutura esperada do template:
        # slide 0 = capa
        # slide 1 = página intermediária em branco para conteúdo
        # último slide = página de fechamento, que não deve ser preenchida
        if len(prs.slides) < 3:
            st.error("O template precisa ter pelo menos 3 slides: capa, página de conteúdo e página de fechamento.")
            st.stop()

        # Usa o slide 1 como primeira página de conteúdo.
        # Ele é limpo e recriado pelo código, mantendo capa e fechamento intactos.
        slide_atual = prs.slides[1]
        limpar_slide(slide_atual)
        adicionar_logo_rodape(slide_atual)
        y_atual = TOPO_CONTEUDO

        for titulo, texto, grafico in secoes:
            if not titulo.strip() and not texto.strip() and grafico is None:
                continue

            slide_atual, y_atual = adicionar_secao_no_relatorio(
                prs=prs,
                slide_atual=slide_atual,
                titulo=titulo,
                texto=texto,
                y_atual=y_atual,
                grafico=grafico
            )

        output = BytesIO()
        prs.save(output)
        output.seek(0)

        st.download_button(
            label="Baixar Morning Call em PowerPoint",
            data=output,
            file_name="morning_call.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )