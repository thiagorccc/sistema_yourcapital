import io
import math
import os
import re
import streamlit as st
import numpy as np
import pandas as pd
import yfinance as yf
import plotly.graph_objects as go
import plotly.figure_factory as ff
import empyrical
import time
from PIL import Image
from scipy.optimize import minimize

from Assets import symbol_real, symbol_dolar, symbol_euro, symbol_br_indices, load_br_indices, BR_BENCH_TICKERS


symbol_info = symbol_real | symbol_dolar | symbol_euro | symbol_br_indices


def _load_cdi_bcb(start, end):
    """Carrega CDI via API pública do BCB/SGS com chunking de 9 anos (limite da API)."""
    import requests
    fmt = "%d/%m/%Y"
    start_dt, end_dt = pd.to_datetime(start), pd.to_datetime(end)
    rows, cur = [], start_dt
    while cur <= end_dt:
        nxt = min(cur + pd.DateOffset(years=9) - pd.Timedelta(days=1), end_dt)
        resp = requests.get(
            "https://api.bcb.gov.br/dados/serie/bcdata.sgs.12/dados",
            params={"formato": "json",
                    "dataInicial": cur.strftime(fmt),
                    "dataFinal": nxt.strftime(fmt)},
            timeout=30,
        )
        resp.raise_for_status()
        rows.extend(resp.json())
        cur = nxt + pd.Timedelta(days=1)
    dates = pd.to_datetime([d["data"] for d in rows], dayfirst=True)
    values = [float(d["valor"]) / 100.0 for d in rows]
    s = pd.Series(values, index=dates, name="CDI").sort_index()
    return s[~s.index.duplicated(keep="last")]


def _compute_synthetic_returns(ref_ticker, start, end):
    """Retorna retornos diários do ticker de referência de um ativo sintético.
    Suporta índices BR (via _get_br_slice) e qualquer ticker do yfinance."""
    if ref_ticker in symbol_br_indices:
        br = _get_br_slice(start=str(start), end=str(end))
        if ref_ticker in br.columns:
            return br[ref_ticker].dropna()
        if ref_ticker == "CDI":
            try:
                s = _load_cdi_bcb(start, end)
                return s.dropna() if not s.empty else None
            except Exception:
                pass
        return None
    try:
        df = yf.download(ref_ticker, start=start, end=end, progress=False, auto_adjust=True)
        if df.empty:
            return None
        price = df["Close"].squeeze() if isinstance(df, pd.DataFrame) else df
        return pd.Series(price).pct_change().dropna()
    except Exception:
        return None


@st.cache_resource(show_spinner="Carregando índices BR (CDI, IMA-B, IHFA...)...")
def _get_br_full():
    """Carrega e mantém os índices BR em cache para toda a sessão do servidor."""
    return load_br_indices()


def _get_br_slice(start=None, end=None):
    full = _get_br_full()
    if full.empty:
        return full
    result = full.copy()
    if start is not None:
        result = result[result.index >= pd.to_datetime(start)]
    if end is not None:
        result = result[result.index <= pd.to_datetime(end)]
    return result

benchmarks = {
    "Ibovespa (Brazil)": "^BVSP",
    "S&P 500 (SPY)": "SPY",
    "S&P 500 Index (^GSPC)": "^GSPC",
    "Dow Jones (^DJI)": "^DJI",
    "Nasdaq 100 (^NDX)": "^NDX",
    "Russell 2000 (^RUT)": "^RUT",
    "EURO STOXX 50": "^STOXX50E",
    "STOXX Europe 600": "^STOXX",
    "DAX (Germany)": "^GDAXI",
    "CAC 40 (France)": "^FCHI",
    "FTSE MIB (Italy)": "FTSEMIB.MI",
    "IBEX 35 (Spain)": "^IBEX",
    "FTSE 100 (UK)": "^FTSE",
    "Euro Corporate Bonds (ETF - IEAC.L)": "IEAC.L",
    "U.S. Treasury Bill ETF (Cash Equivalent)": "BIL",
    "CDI (Taxa DI)": "CDI",
    "IMA-B (NTN-B Total)": "IMA-B",
    "IHFA (Hedge Funds)": "IHFA",
}


def safe_download(symbol, start, end, max_retries=10, wait=2):
    if symbol in symbol_br_indices:
        return pd.Series(dtype=float)
    for _ in range(max_retries):
        try:
            df = yf.download(symbol, start=start, end=end, progress=False, auto_adjust=True)
            if isinstance(df.columns, pd.MultiIndex):
                df.columns = df.columns.get_level_values(0)
            if "Close" in df and not df["Close"].dropna().empty:
                close = df["Close"]
                if isinstance(close, pd.DataFrame):
                    close = close.squeeze()
                return pd.Series(close).dropna()
        except Exception:
            pass
        time.sleep(wait)
    return pd.Series(dtype=float)


def _ajustar_cambio(returns, currency_choice, today):
    if currency_choice == "Real":
        for symbol in returns.columns:
            if symbol in symbol_dolar:
                dollar_brl = safe_download("USDBRL=X", start="2000-01-01", end=today)
                aligned = np.log(dollar_brl / dollar_brl.shift(1)).dropna()
                aligned = aligned.reindex(returns.index).dropna()
                returns = returns.loc[aligned.index]
                returns[symbol] = (1 + returns[symbol]) * (1 + aligned) - 1
            elif symbol in symbol_euro:
                euro_brl = safe_download("EURBRL=X", start="2000-01-01", end=today)
                aligned = np.log(euro_brl / euro_brl.shift(1)).dropna()
                aligned = aligned.reindex(returns.index).dropna()
                returns = returns.loc[aligned.index]
                returns[symbol] = (1 + returns[symbol]) * (1 + aligned) - 1
    elif currency_choice == "Dollar":
        for symbol in returns.columns:
            if symbol in symbol_real:
                dollar_brl = safe_download("USDBRL=X", start="2000-01-01", end=today)
                aligned = -np.log(dollar_brl / dollar_brl.shift(1)).dropna()
                aligned = aligned.reindex(returns.index).dropna()
                returns = returns.loc[aligned.index]
                returns[symbol] = (1 + returns[symbol]) * (1 + aligned) - 1
            elif symbol in symbol_euro:
                usd_eur = safe_download("USDEUR=X", start="2000-01-01", end=today)
                aligned = np.log(usd_eur / usd_eur.shift(1)).dropna()
                aligned = aligned.reindex(returns.index).dropna()
                returns = returns.loc[aligned.index]
                returns[symbol] = (1 + returns[symbol]) * (1 + aligned) - 1
    elif currency_choice == "Euro":
        for symbol in returns.columns:
            if symbol in symbol_real:
                euro_brl = safe_download("EURBRL=X", start="2000-01-01", end=today)
                aligned = -np.log(euro_brl / euro_brl.shift(1)).dropna()
                aligned = aligned.reindex(returns.index).dropna()
                returns = returns.loc[aligned.index]
                returns[symbol] = (1 + returns[symbol]) * (1 + aligned) - 1
            elif symbol in symbol_dolar:
                usd_eur = safe_download("USDEUR=X", start="2000-01-01", end=today)
                aligned = -np.log(usd_eur / usd_eur.shift(1)).dropna()
                aligned = aligned.reindex(returns.index).dropna()
                returns = returns.loc[aligned.index]
                returns[symbol] = (1 + returns[symbol]) * (1 + aligned) - 1
    return returns


def _ajustar_cambio_benchmark(benchmark_returns, benchmark_ticker, currency_choice, today):
    if benchmark_returns.empty or benchmark_ticker is None:
        return benchmark_returns
    if currency_choice == "Real" and benchmark_ticker in symbol_dolar:
        dollar_brl = safe_download("USDBRL=X", start="2000-01-01", end=today)
        aligned_bm = np.log(dollar_brl / dollar_brl.shift(1)).dropna()
        aligned_bm = aligned_bm.reindex(benchmark_returns.index).dropna()
        benchmark_returns = benchmark_returns.loc[aligned_bm.index]
        benchmark_returns = (1 + benchmark_returns) * (1 + aligned_bm.squeeze()) - 1
    elif currency_choice == "Real" and benchmark_ticker in symbol_euro:
        euro_brl = safe_download("EURBRL=X", start="2000-01-01", end=today)
        aligned_bm = np.log(euro_brl / euro_brl.shift(1)).dropna()
        aligned_bm = aligned_bm.reindex(benchmark_returns.index).dropna()
        benchmark_returns = benchmark_returns.loc[aligned_bm.index]
        benchmark_returns = (1 + benchmark_returns) * (1 + aligned_bm.squeeze()) - 1
    elif currency_choice == "Dollar" and benchmark_ticker in symbol_real:
        dollar_brl = safe_download("USDBRL=X", start="2000-01-01", end=today)
        aligned_bm = -np.log(dollar_brl / dollar_brl.shift(1)).dropna()
        aligned_bm = aligned_bm.reindex(benchmark_returns.index).dropna()
        benchmark_returns = benchmark_returns.loc[aligned_bm.index]
        benchmark_returns = (1 + benchmark_returns) * (1 + aligned_bm.squeeze()) - 1
    elif currency_choice == "Dollar" and benchmark_ticker in symbol_euro:
        usd_eur = safe_download("USDEUR=X", start="2000-01-01", end=today)
        aligned_bm = np.log(usd_eur / usd_eur.shift(1)).dropna()
        aligned_bm = aligned_bm.reindex(benchmark_returns.index).dropna()
        benchmark_returns = benchmark_returns.loc[aligned_bm.index]
        benchmark_returns = (1 + benchmark_returns) * (1 + aligned_bm.squeeze()) - 1
    elif currency_choice == "Euro" and benchmark_ticker in symbol_real:
        euro_brl = safe_download("EURBRL=X", start="2000-01-01", end=today)
        aligned_bm = -np.log(euro_brl / euro_brl.shift(1)).dropna()
        aligned_bm = aligned_bm.reindex(benchmark_returns.index).dropna()
        benchmark_returns = benchmark_returns.loc[aligned_bm.index]
        benchmark_returns = (1 + benchmark_returns) * (1 + aligned_bm.squeeze()) - 1
    elif currency_choice == "Euro" and benchmark_ticker in symbol_dolar:
        usd_eur = safe_download("USDEUR=X", start="2000-01-01", end=today)
        aligned_bm = -np.log(usd_eur / usd_eur.shift(1)).dropna()
        aligned_bm = aligned_bm.reindex(benchmark_returns.index).dropna()
        benchmark_returns = benchmark_returns.loc[aligned_bm.index]
        benchmark_returns = (1 + benchmark_returns) * (1 + aligned_bm.squeeze()) - 1
    return benchmark_returns


def _render_synth_expander(key_prefix, ticker_options, ticker_lookup):
    """Expander compacto para criar/remover ativos sintéticos, exibido dentro de cada tab."""
    synth = st.session_state.setdefault("synthetic_assets", {})
    label = f"Ativos Sintéticos ({len(synth)} criado{'s' if len(synth) != 1 else ''})" if synth else "Ativos Sintéticos"
    # Filtra sintéticos já criados para não aparecerem como referência deles mesmos
    ref_options = [o for o in ticker_options if not any(
        o.endswith(f"[{sn}]") for sn in synth
    )]
    with st.expander(label, expanded=False):
        if synth:
            for sname, info in list(synth.items()):
                c1, c2, c3, c4 = st.columns([2, 2, 1, 1])
                c1.write(f"**{sname}**")
                c2.write(info["ref_ticker"])
                c3.write(f"β = {info['beta']:.2g}")
                if c4.button("✕", key=f"rm_{key_prefix}_{sname}", help=f"Remover {sname}"):
                    del st.session_state["synthetic_assets"][sname]
                    st.rerun()
            st.divider()
        with st.form(f"form_synth_{key_prefix}", clear_on_submit=True):
            fc1, fc2, fc3 = st.columns(3)
            new_name    = fc1.text_input("Nome (ticker)", placeholder="Ex: XPTO")
            ref_display = fc2.selectbox("Referência", options=ref_options)
            new_beta    = fc3.number_input("Beta", min_value=0.0, max_value=20.0, value=1.0, step=0.01, format="%.2f")
            if st.form_submit_button("Adicionar"):
                nc = new_name.strip().upper()
                rc = ticker_lookup.get(ref_display, ref_display)
                if nc and rc:
                    st.session_state["synthetic_assets"][nc] = {"ref_ticker": rc, "beta": float(new_beta)}
                    st.rerun()
                else:
                    st.error("Preencha o nome do ativo.")


def _render_constraints_expander(tickers, key_prefix):
    """Expander para restrições de alocação: global, por ativo e por grupo.
    Retorna (g_min, g_max, per_asset_bounds_dict, group_constraints_list)."""
    g_min, g_max, per_asset, groups = 0.0, 1.0, {}, []
    with st.expander("Restrições de Alocação", expanded=False):
        st.markdown("##### Limites Globais")
        st.caption("Aplicados a todos os ativos sem restrição individual definida abaixo.")
        cg1, cg2 = st.columns(2)

        # Restore number_input values from backup keys (lost when "Voltar" hides widgets)
        _gmin_key = f"alloc_gmin_{key_prefix}"
        _gmax_key = f"alloc_gmax_{key_prefix}"
        if _gmin_key not in st.session_state:
            st.session_state[_gmin_key] = st.session_state.get(f"{_gmin_key}_bak", 0.0)
        if _gmax_key not in st.session_state:
            st.session_state[_gmax_key] = st.session_state.get(f"{_gmax_key}_bak", 100.0)

        g_min_pct = cg1.number_input(
            "Mínimo por ativo (%)", min_value=0.0, max_value=50.0, value=0.0, step=0.5,
            key=_gmin_key,
        )
        g_max_pct = cg2.number_input(
            "Máximo por ativo (%)", min_value=1.0, max_value=100.0, value=100.0, step=0.5,
            key=_gmax_key,
        )
        st.session_state[f"{_gmin_key}_bak"] = g_min_pct
        st.session_state[f"{_gmax_key}_bak"] = g_max_pct
        g_min, g_max = g_min_pct / 100, g_max_pct / 100

        st.markdown("##### Restrições por Ativo")
        st.caption("Adicione linhas para sobrescrever os limites globais em ativos específicos.")
        ticker_opts = [f"{symbol_info.get(t, t)} [{t}]" for t in tickers]
        ticker_map  = {f"{symbol_info.get(t, t)} [{t}]": t for t in tickers}
        empty_override = pd.DataFrame({
            "Ativo":   pd.Series(dtype="str"),
            "Mín (%)": pd.Series(dtype="float"),
            "Máx (%)": pd.Series(dtype="float"),
        })
        _override_snap = f"alloc_override_snap_{key_prefix}"
        _override_key  = f"alloc_override_{key_prefix}"
        _override_seed = f"alloc_override_seed_{key_prefix}"
        # A semente só pode mudar de valor no render em que a chave do widget é
        # (re)criada — passar dados diferentes em qualquer render subsequente,
        # mesmo com a chave já existindo, faz o Streamlit perder a correlação
        # com o que está sendo editado (mesma causa do reset da tabela de
        # ativos após "Voltar", já corrigido em _portfolio_editor).
        if _override_key not in st.session_state:
            st.session_state[_override_seed] = (
                st.session_state[_override_snap].copy()
                if _override_snap in st.session_state
                else empty_override
            )
        override_df = st.data_editor(
            st.session_state[_override_seed],
            num_rows="dynamic",
            column_config={
                "Ativo":   st.column_config.SelectboxColumn("Ativo", options=ticker_opts, required=True),
                "Mín (%)": st.column_config.NumberColumn("Mín (%)", min_value=0.0, max_value=100.0, step=0.1, format="%.2f"),
                "Máx (%)": st.column_config.NumberColumn("Máx (%)", min_value=0.0, max_value=100.0, step=0.1, format="%.2f"),
            },
            hide_index=True,
            key=_override_key,
        )
        st.session_state[_override_snap] = override_df.reset_index(drop=True).copy()

        st.markdown("##### Restrições por Grupo de Ativos")
        st.caption("Limites para o peso total de um conjunto de ativos (ex: máx 60% em renda variável). Use os tickers exatos, separados por vírgula.")
        empty_groups = pd.DataFrame({
            "Nome do Grupo":               pd.Series(dtype="str"),
            "Ativos (tickers, vírgula)":   pd.Series(dtype="str"),
            "Mín (%)":                     pd.Series(dtype="float"),
            "Máx (%)":                     pd.Series(dtype="float"),
        })
        _groups_snap = f"alloc_groups_snap_{key_prefix}"
        _groups_key  = f"alloc_groups_{key_prefix}"
        _groups_seed = f"alloc_groups_seed_{key_prefix}"
        # Mesma lógica de semente congelada usada acima para o override_df.
        if _groups_key not in st.session_state:
            st.session_state[_groups_seed] = (
                st.session_state[_groups_snap].copy()
                if _groups_snap in st.session_state
                else empty_groups
            )
        group_df = st.data_editor(
            st.session_state[_groups_seed],
            num_rows="dynamic",
            column_config={
                "Nome do Grupo":             st.column_config.TextColumn("Nome do Grupo"),
                "Ativos (tickers, vírgula)": st.column_config.TextColumn("Ativos"),
                "Mín (%)": st.column_config.NumberColumn("Mín (%)", min_value=0.0, max_value=100.0, step=0.5, format="%.1f"),
                "Máx (%)": st.column_config.NumberColumn("Máx (%)", min_value=0.0, max_value=100.0, step=0.5, format="%.1f"),
            },
            hide_index=True,
            key=_groups_key,
        )
        st.session_state[_groups_snap] = group_df.reset_index(drop=True).copy()

    # Parse per-asset overrides
    for _, row in override_df.dropna(subset=["Ativo"]).iterrows():
        raw_t = ticker_map.get(str(row["Ativo"]), str(row["Ativo"]))
        mn = float(row["Mín (%)"]) / 100 if pd.notna(row.get("Mín (%)")) else g_min
        mx = float(row["Máx (%)"]) / 100 if pd.notna(row.get("Máx (%)")) else g_max
        per_asset[raw_t] = (mn, mx)

    # Parse group constraints
    for _, row in group_df.iterrows():
        raw_assets = str(row.get("Ativos (tickers, vírgula)", "")).strip()
        if not raw_assets or raw_assets == "nan":
            continue
        grp_tickers = [t.strip().upper() for t in raw_assets.split(",") if t.strip()]
        mn = float(row["Mín (%)"]) / 100 if pd.notna(row.get("Mín (%)")) else 0.0
        mx = float(row["Máx (%)"]) / 100 if pd.notna(row.get("Máx (%)")) else 1.0
        groups.append({"name": str(row.get("Nome do Grupo", "")), "tickers": grp_tickers, "min": mn, "max": mx})

    return g_min, g_max, per_asset, groups


def _build_opt_constraints(ordered_tickers, groups):
    """Constrói lista de restrições scipy.minimize para grupos de ativos."""
    cons = []
    for gc in groups:
        idx = [i for i, t in enumerate(ordered_tickers) if t in gc["tickers"]]
        if not idx:
            continue
        mn, mx = gc["min"], gc["max"]
        if mn > 0:
            cons.append({"type": "ineq", "fun": lambda w, i=idx, m=mn: np.array(w)[i].sum() - m})
        if mx < 1.0:
            cons.append({"type": "ineq", "fun": lambda w, i=idx, m=mx: m - np.array(w)[i].sum()})
    return cons


_CURRENCY_SYMBOLS = {"Real": "R$", "Dollar": "US$", "Euro": "€"}
_MODE_PCT = "Percentual (%)"
_MODE_VAL = "Valor financeiro"


def _portfolio_editor(default_symbols, default_weights, ticker_options, ticker_lookup, key_prefix, currency_choice="Real"):
    editor_key       = f"{key_prefix}_editor"
    snapshot_key     = f"{key_prefix}_snapshot"
    seed_key         = f"{key_prefix}_seed"
    editor_key_val   = f"{key_prefix}_editor_val"
    snapshot_key_val = f"{key_prefix}_snapshot_val"
    seed_key_val     = f"{key_prefix}_seed_val"
    mode_key         = f"{key_prefix}_mode"
    mode_bak_key     = f"{mode_key}_bak"
    total_key        = f"{key_prefix}_total_value"
    total_bak_key    = f"{total_key}_bak"
    sym = _CURRENCY_SYMBOLS.get(currency_choice, "R$")

    # Modo de preenchimento (% ou valor financeiro). Restaurado via chave
    # "_bak" — mesmo idioma já usado em _render_constraints_expander para
    # widgets simples que somem do session_state quando "Voltar" esconde o
    # bloco. Cada modo mantém seus próprios dados de forma independente —
    # não há conversão automática ao alternar entre eles.
    if mode_key not in st.session_state:
        st.session_state[mode_key] = st.session_state.get(mode_bak_key, _MODE_PCT)
    mode = st.radio("Preencher alocação por:", [_MODE_PCT, _MODE_VAL], key=mode_key, horizontal=True)
    st.session_state[mode_bak_key] = mode

    if mode == _MODE_VAL:
        if total_key not in st.session_state:
            st.session_state[total_key] = st.session_state.get(total_bak_key, 1_000_000.0)
        total_value = st.number_input(
            f"Valor total da carteira ({sym})",
            min_value=0.0, value=1_000_000.0, step=10_000.0, format="%.2f",
            key=total_key,
        )
        st.session_state[total_bak_key] = total_value
    else:
        total_value = st.session_state.get(total_bak_key, 1_000_000.0)

    if mode == _MODE_PCT:
        # `default_data` só pode mudar de valor no render em que a chave do
        # widget é (re)criada — Streamlit rastreia as edições do usuário como
        # deltas relativos ao `data=` daquele momento, então se o `data=`
        # passado mudar em qualquer render subsequente (mesmo com a chave já
        # existindo), o widget perde a correlação com o que está editando e
        # tanto duplica edições quanto, silenciosamente, volta a exibir o novo
        # `data=` para células nunca tocadas nesta "vida" do widget. Por isso
        # a semente é calculada uma única vez (ao criar/recriar a chave, ex.:
        # após "Voltar") e fica congelada em `seed_key` enquanto o widget
        # existir; o snapshot (sempre atualizado) só é lido para gerar essa
        # semente, nunca passado direto como `data=`.
        if editor_key not in st.session_state:
            if snapshot_key in st.session_state:
                st.session_state[seed_key] = st.session_state[snapshot_key].copy()
            else:
                st.session_state[seed_key] = pd.DataFrame({
                    "Ticker": [f"{symbol_info.get(s, s)} [{s}]" for s in default_symbols],
                    "Weight": [w * 100 for w in default_weights],
                })
        default_data = st.session_state[seed_key]

        display = st.data_editor(
            default_data,
            num_rows="dynamic",
            column_config={
                "Ticker": st.column_config.SelectboxColumn("Ativo", options=ticker_options),
                "Weight": st.column_config.NumberColumn("Peso (%)", min_value=0.0, max_value=100.0, step=0.01, format="%.2f"),
            },
            hide_index=True,
            key=editor_key,
        )
        st.session_state[snapshot_key] = display.reset_index(drop=True).copy()
        portfolio = display.dropna(subset=["Ticker"]).copy()
        portfolio["Ticker"] = portfolio["Ticker"].map(ticker_lookup)
        portfolio["Weight"] = portfolio["Weight"] / 100
        weight_sum = portfolio["Weight"].sum() * 100

    else:  # _MODE_VAL
        # Mesma lógica de semente congelada do ramo percentual — ver comentário acima.
        if editor_key_val not in st.session_state:
            if snapshot_key_val in st.session_state:
                st.session_state[seed_key_val] = st.session_state[snapshot_key_val].copy()
            else:
                st.session_state[seed_key_val] = pd.DataFrame({
                    "Ticker": [f"{symbol_info.get(s, s)} [{s}]" for s in default_symbols],
                    "Valor": [w * total_value for w in default_weights],
                })
        default_data_val = st.session_state[seed_key_val]

        display_val = st.data_editor(
            default_data_val,
            num_rows="dynamic",
            column_config={
                "Ticker": st.column_config.SelectboxColumn("Ativo", options=ticker_options),
                "Valor": st.column_config.NumberColumn(f"Valor ({sym})", min_value=0.0, step=100.0, format="%.2f"),
            },
            hide_index=True,
            key=editor_key_val,
        )
        st.session_state[snapshot_key_val] = display_val.reset_index(drop=True).copy()
        portfolio = display_val.dropna(subset=["Ticker"]).copy()
        portfolio["Ticker"] = portfolio["Ticker"].map(ticker_lookup)
        if total_value > 0:
            portfolio["Weight"] = portfolio["Valor"] / total_value
        else:
            portfolio["Weight"] = 0.0
            st.warning("Informe um valor total da carteira maior que zero para calcular os pesos.")
        portfolio = portfolio.drop(columns=["Valor"])
        weight_sum = portfolio["Weight"].sum() * 100

        alocado = display_val["Valor"].fillna(0).sum()
        # Escapa "$" — com dois símbolos de moeda na mesma string (ex.: "R$ ... de R$ ..."),
        # o st.caption interpretaria o trecho entre eles como LaTeX (delimitador $...$).
        _valor_alocado_txt = f"Valor alocado: {_brl(alocado, symbol=sym)} de {_brl(total_value, symbol=sym)}"
        st.caption(_valor_alocado_txt.replace("$", "\\$"))

    st.markdown(f"**Alocação Total: {weight_sum:.2f}%**")
    return portfolio, weight_sum


def _tabela_rentabilidade_mensal(portfolios_dict, benchmark_series=None, benchmark_label=None):
    """Tabela de rentabilidade mensal no estilo factsheet.
    Para CDI: exibe linha '% CDI' abaixo de cada carteira.
    Para outro benchmark: exibe linha do benchmark abaixo de todas as carteiras."""
    MESES = ["Jan", "Fev", "Mar", "Abr", "Mai", "Jun", "Jul", "Ago", "Set", "Out", "Nov", "Dez"]
    is_cdi = benchmark_label is not None and "CDI" in str(benchmark_label)

    def _norm(s):
        s = s.copy()
        if getattr(s.index, "tz", None) is not None:
            s.index = s.index.tz_localize(None)
        s.index = s.index.normalize()
        return s.dropna().sort_index()

    def _to_monthly(s):
        try:
            return (1 + _norm(s)).resample("ME").prod() - 1
        except ValueError:
            return (1 + _norm(s)).resample("M").prod() - 1

    mp = {lbl: _to_monthly(s) for lbl, s in portfolios_dict.items() if not s.empty}
    mb = _to_monthly(benchmark_series) if (benchmark_series is not None and not benchmark_series.empty) else None

    if not mp:
        return

    all_years = sorted({yr for s in mp.values() for yr in s.index.year}, reverse=True)

    # Data de início comum: usa a mais tardia entre carteiras e benchmark,
    # evitando que CDI carregado desde 2001 distorça o acumulado de uma
    # carteira com ativos disponíveis apenas a partir de 2010, por exemplo.
    all_starts = [s.index.min() for s in mp.values()]
    if mb is not None:
        all_starts.append(mb.index.min())
    common_start = max(all_starts)

    def _get_vals(s_m, year):
        yd = s_m[s_m.index.year == year]
        return [
            yd[yd.index.month == m].iloc[0] if len(yd[yd.index.month == m]) > 0 else np.nan
            for m in range(1, 13)
        ]

    def _ytd(vals):
        valid = [v for v in vals if not pd.isna(v)]
        return (np.prod([1 + v for v in valid]) - 1) if valid else np.nan

    def _acc(s_m, year):
        d = s_m[(s_m.index >= common_start) & (s_m.index.year <= year)].dropna()
        return ((1 + d).prod() - 1) if not d.empty else np.nan

    def _td(v, small=False, bold=False, is_rel=False, border_top=False):
        bt = "border-top:2px solid #e2e8f0;" if border_top else ""
        sz = "11px" if small else "13px"
        fw = "bold" if bold else "normal"
        if pd.isna(v):
            return f'<td style="text-align:right;padding:3px 8px;font-size:{sz};font-weight:{fw};{bt}">—</td>'
        if is_rel:
            color = "#dc2626" if v < 100 else "#059669"
            return f'<td style="text-align:right;padding:3px 8px;font-size:{sz};color:{color};font-weight:{fw};{bt}">{v:.2f}%</td>'
        else:
            color = "#dc2626" if v < 0 else "inherit"
            return f'<td style="text-align:right;padding:3px 8px;font-size:{sz};color:{color};font-weight:{fw};{bt}">{v*100:.2f}%</td>'

    rows_per_year = len(mp) * (2 if is_cdi else 1) + (0 if is_cdi or mb is None else 1)

    body = []
    for year in all_years:
        first = True
        for lbl, s_m in mp.items():
            pv = _get_vals(s_m, year)
            ytd_v = _ytd(pv)
            acc_v = _acc(s_m, year)

            year_td = ""
            if first:
                year_td = (
                    f'<td rowspan="{rows_per_year}" style="font-weight:bold;padding:4px 10px;'
                    f'border-top:2px solid #e2e8f0;vertical-align:middle;font-size:14px;text-align:left;">{year}</td>'
                )
                first = False

            name_td = f'<td style="padding:3px 8px;border-top:2px solid #e2e8f0;white-space:nowrap;text-align:left;">{lbl}</td>'
            month_tds = "".join(_td(v, border_top=True) for v in pv)
            body.append(
                f'<tr>{year_td}{name_td}{month_tds}'
                f'{_td(ytd_v, bold=True, border_top=True)}{_td(acc_v, bold=True, border_top=True)}</tr>'
            )

            if is_cdi and mb is not None:
                bm_vals = _get_vals(mb, year)
                bm_ytd = _ytd(bm_vals)
                acc_bm = _acc(mb, year)
                rel_tds = ""
                for p, b in zip(pv, bm_vals):
                    rel = (p / b * 100) if (not pd.isna(p) and not pd.isna(b) and b > 0) else np.nan
                    rel_tds += _td(rel, small=True, is_rel=True)
                ytd_rel = (ytd_v / bm_ytd * 100) if (not pd.isna(ytd_v) and not pd.isna(bm_ytd) and bm_ytd > 0) else np.nan
                acc_rel = (acc_v / acc_bm * 100) if (not pd.isna(acc_v) and not pd.isna(acc_bm) and acc_bm > 0) else np.nan
                pct_td = '<td style="padding:2px 8px;color:#64748b;font-size:11px;text-align:left;">% CDI</td>'
                body.append(
                    f'<tr>{pct_td}{rel_tds}'
                    f'{_td(ytd_rel, small=True, bold=True, is_rel=True)}{_td(acc_rel, small=True, bold=True, is_rel=True)}</tr>'
                )

        if not is_cdi and mb is not None:
            bm_vals = _get_vals(mb, year)
            bm_ytd = _ytd(bm_vals)
            bm_acc = _acc(mb, year)
            bm_name = benchmark_label or "Benchmark"
            bm_name_td = (
                f'<td style="padding:3px 8px;color:#64748b;font-size:12px;'
                f'font-style:italic;text-align:left;">{bm_name}</td>'
            )
            bm_month_tds = "".join(_td(v, small=True) for v in bm_vals)
            body.append(
                f'<tr>{bm_name_td}{bm_month_tds}'
                f'{_td(bm_ytd, small=True, bold=True)}{_td(bm_acc, small=True, bold=True)}</tr>'
            )

    hcols = ["ANO", ""] + MESES + ["No ano", "Acumulado"]
    header = "".join(
        f'<th style="padding:6px 8px;text-align:{"left" if i < 2 else "right"};'
        f'white-space:nowrap;background:#1e293b;color:white;">{c}</th>'
        for i, c in enumerate(hcols)
    )
    html = (
        '<div style="overflow-x:auto;">'
        '<table style="width:100%;border-collapse:collapse;font-family:sans-serif;">'
        f'<thead><tr>{header}</tr></thead>'
        f'<tbody>{"".join(body)}</tbody>'
        '</table></div>'
    )
    st.markdown("### Rentabilidade Mensal")
    st.markdown(html, unsafe_allow_html=True)


def _get_rf_rate(currency, start, end, index):
    """Retorna taxa diária livre de risco média no período.
    Real → CDI  |  Dollar → BIL  |  Euro → EXX6.DE (Govs curtos zona do euro)"""
    today_str = str(pd.Timestamp.today().date())
    if currency == "Real":
        br = _get_br_slice(start=str(start), end=str(end))
        if "CDI" in br.columns:
            cdi = br["CDI"].dropna()
            val = cdi.reindex(index).dropna().mean()
            return float(val) if pd.notna(val) else 0.0
        try:
            cdi = _load_cdi_bcb(start, end)
            if not cdi.empty:
                val = cdi.reindex(index).dropna().mean()
                return float(val) if pd.notna(val) else 0.0
        except Exception:
            pass
        return 0.0
    elif currency == "Euro":
        rf_data = safe_download("EXX6.DE", start="2000-01-01", end=today_str)
        if not rf_data.empty:
            rf_ret = rf_data.pct_change().dropna()
            sl = rf_ret[str(start):str(end)].reindex(index).dropna()
            return float(sl.mean()) if not sl.empty else 0.0
        return 0.0
    else:  # Dollar
        rf_data = safe_download("BIL", start="2000-01-01", end=today_str)
        if not rf_data.empty:
            rf_ret = rf_data.pct_change().dropna()
            sl = rf_ret[str(start):str(end)].reindex(index).dropna()
            return float(sl.mean()) if not sl.empty else 0.0
        return 0.0


def _calcular_metricas(series, rf, label):
    m = float(series.mean())
    s = float(series.std())
    sharpe = ((m - rf) / s) * np.sqrt(252) if s > 0 else np.nan
    downside = np.minimum(series - rf, 0)
    dd_dev = float(np.sqrt(np.mean(np.square(downside))))
    sortino = ((m - rf) / dd_dev) * np.sqrt(252) if dd_dev > 0 else np.nan
    return {
        "label": label,
        "retorno_anual": empyrical.annual_return(series),
        "volatilidade": empyrical.annual_volatility(series),
        "sharpe": sharpe,
        "sortino": sortino,
        "max_drawdown": empyrical.max_drawdown(series),
        "var_1pct": empyrical.value_at_risk(series, cutoff=0.01),
    }


def _exibir_metricas(metricas_list):
    # (nome_exibição, chave, formato, maior_é_melhor)
    ROWS = [
        ("Retorno Anual",    "retorno_anual", ".2%", True),
        ("Volatilidade",     "volatilidade",  ".2%", False),
        ("Índice de Sharpe", "sharpe",        ".2f", True),
        ("Índice de Sortino","sortino",        ".2f", True),
        ("Drawdown Máximo",  "max_drawdown",  ".2%", True),
        ("VaR 1%",           "var_1pct",      ".2%", True),
    ]

    raw, fmt_str, directions = {}, {}, {}
    for nome, chave, fmt, higher_is_better in ROWS:
        raw[nome]     = {m["label"]: (m[chave] if not pd.isna(m[chave]) else np.nan) for m in metricas_list}
        fmt_str[nome] = {m["label"]: (f"{m[chave]:{fmt}}" if not pd.isna(m[chave]) else "N/D") for m in metricas_list}
        directions[nome] = higher_is_better

    df_raw  = pd.DataFrame(raw).T
    df_disp = pd.DataFrame(fmt_str).T
    df_disp.index.name = "Métrica"

    def _highlight(df):
        styles = pd.DataFrame("", index=df.index, columns=df.columns)
        if len(metricas_list) < 2:
            return styles
        for nome in df.index:
            valid = df_raw.loc[nome].dropna()
            if len(valid) < 2:
                continue
            higher = directions[nome]
            best  = valid.idxmax() if higher else valid.idxmin()
            worst = valid.idxmin() if higher else valid.idxmax()
            styles.loc[nome, best]  = "background-color: #d1fae5; color: #065f46; font-weight: bold"
            if best != worst:
                styles.loc[nome, worst] = "background-color: #fee2e2; color: #991b1b"
        return styles

    st.dataframe(
        df_disp.style.apply(_highlight, axis=None),
        use_container_width=True,
    )


def _graficos_comparacao(series_dict):
    # Normaliza timezone e remove componente de hora para garantir compatibilidade
    # entre séries do yfinance (pode ser tz-aware) e do BCB/ANBIMA (tz-naive).
    def _normalizar(s):
        s = s.copy()
        if getattr(s.index, "tz", None) is not None:
            s.index = s.index.tz_localize(None)
        s.index = s.index.normalize()
        return s.sort_index()

    normed = {lbl: _normalizar(s) for lbl, s in series_dict.items() if not s.empty}
    if not normed:
        st.warning("Sem dados suficientes para exibir os gráficos.")
        return

    # Índice de referência = primeira série (geralmente a carteira principal)
    ref_idx = next(iter(normed.values())).index

    # Alinha todas as séries ao índice de referência.
    # ffill preenche feriados de calendários diferentes (ex: CDI em feriado BR
    # vs dia útil americano), garantindo que o benchmark apareça no gráfico.
    aligned = {}
    for lbl, s in normed.items():
        s_al = s.reindex(ref_idx).ffill().dropna()
        if not s_al.empty:
            aligned[lbl] = s_al

    if not aligned:
        st.warning("Sem dados suficientes para exibir os gráficos.")
        return

    # Garante que todas as séries começam na mesma data — evita descasamento
    # visual quando o CDI/benchmark começa 1-2 dias depois da carteira.
    common_start = max(s.index.min() for s in aligned.values())
    aligned = {lbl: s[s.index >= common_start] for lbl, s in aligned.items()}

    st.markdown("### Retorno Acumulado")
    fig = go.Figure()
    for lbl, s in aligned.items():
        fig.add_trace(go.Scatter(x=s.index, y=(1 + s).cumprod(), mode="lines", name=lbl))
    fig.update_layout(title="Comparação de Retorno Acumulado", xaxis_title="Data", yaxis_title="Valor da Carteira")
    st.plotly_chart(fig, use_container_width=True)

    st.markdown("### Volatilidade Móvel (30 Dias)")
    fig_vol = go.Figure()
    for lbl, s in aligned.items():
        fig_vol.add_trace(go.Scatter(x=s.index, y=s.rolling(30).std() * np.sqrt(252), mode="lines", name=lbl))
    fig_vol.update_layout(title="Comparação de Volatilidade Móvel", xaxis_title="Data", yaxis_title="Volatilidade")
    st.plotly_chart(fig_vol, use_container_width=True)

    st.markdown("### Drawdown da Carteira")
    fig_dd = go.Figure()
    for lbl, s in aligned.items():
        cum = (1 + s).cumprod()
        fig_dd.add_trace(go.Scatter(x=cum.index, y=(cum / cum.cummax()) - 1, mode="lines", name=lbl))
    fig_dd.update_layout(title="Comparação de Drawdown", xaxis_title="Data", yaxis_title="Drawdown")
    st.plotly_chart(fig_dd, use_container_width=True)


def _matriz_correlacao(returns_df):
    if len(returns_df.columns) < 2:
        return
    st.markdown("### Matriz de Correlação dos Ativos")
    corr = returns_df.corr()
    annot = corr.round(2).astype(str).values.tolist()
    fig_corr = ff.create_annotated_heatmap(
        z=corr.values,
        x=corr.columns.tolist(),
        y=corr.index.tolist(),
        annotation_text=annot,
        colorscale="RdBu",
        showscale=True,
        reversescale=True,
        zmin=-1,
        zmax=1,
    )
    st.plotly_chart(fig_corr, use_container_width=True)


def _fig_to_png_bytes(fig, width=900, height=600):
    return fig.to_image(format="png", width=width, height=height, scale=2)


def _align_series_for_report(series_dict):
    def _norm(s):
        s = s.copy()
        if getattr(s.index, "tz", None) is not None:
            s.index = s.index.tz_localize(None)
        s.index = s.index.normalize()
        return s.sort_index()
    normed = {lbl: _norm(s) for lbl, s in series_dict.items() if not s.empty}
    if not normed:
        return {}
    ref_idx = next(iter(normed.values())).index
    aligned = {}
    for lbl, s in normed.items():
        s_al = s.reindex(ref_idx).ffill().dropna()
        if not s_al.empty:
            aligned[lbl] = s_al
    return aligned


_REPORT_PALETTE = ["#60A5FA", "#1E3A8A", "#DC2626", "#059669", "#7C3AED", "#0891B2"]
_YC_NAVY = "#001441"


def _brl(value, decimals=2, symbol="R$"):
    """Format a number in Brazilian-style thousands/decimals: R$ 1.234.567,89"""
    s = f"{abs(value):,.{decimals}f}"
    s = s.replace(",", "X").replace(".", ",").replace("X", ".")
    sign = "-" if value < 0 else ""
    return f"{symbol} {sign}{s}"


# ---------------------------------------------------------------------------
# Textos pré-prontos para os slides do relatório PPTX
# ---------------------------------------------------------------------------
_TEXTOS_PADRAO = {
    "fronteira": {
        "Explicação": (
            "A *fronteira eficiente* é o conjunto de carteiras que oferecem o máximo retorno esperado "
            "para cada nível de risco aceitável, ou equivalentemente, o mínimo risco para cada nível "
            "de retorno desejado. Construída através da teoria moderna de portfolio (Média-Variância), "
            "ela representa as combinações ótimas de ativos que eliminam o risco não-sistemático através "
            "da diversificação. Qualquer carteira localizada abaixo da fronteira é considerada *ineficiente*, "
            "pois oferece menor retorno para o mesmo risco ou maior risco para o mesmo retorno."
        ),
        "Análise do Gráfico": (
            "O gráfico evidencia a *ineficiência da carteira atual*, que se situa abaixo da fronteira "
            "eficiente, oferecendo retorno de {atual_ret} com risco de {atual_risco}. "
            "A análise comparativa apresenta duas alternativas otimizadas: a carteira \"Mesmo Risco\" "
            "mantém o risco em {mesmo_risco_risco} enquanto eleva o retorno para {mesmo_risco_ret}, "
            "e a carteira \"Mesmo Retorno\" oferece o mesmo retorno da atual ({atual_ret}) com risco "
            "reduzido a {mesmo_ret_risco}. *A estratégia recomendada é a implementação da carteira "
            "\"Mesmo Risco\"*, que permite capturar *{ganho_ret} pontos percentuais adicionais de retorno "
            "anual* mantendo a volatilidade controlada, consolidando uma alocação otimizada e "
            "matematicamente superior."
        ),
    },
    "metricas": {
        "Explicação": (
            "As métricas de desempenho consolidam a avaliação quantitativa das carteiras analisadas. "
            "O *Retorno Anualizado* mede o crescimento médio gerado ao longo do período, enquanto a "
            "*Volatilidade Anualizada* mensura a dispersão dos retornos como proxy de risco total. "
            "O *Índice de Sharpe* e o *Índice de Sortino* quantificam o retorno por unidade de risco total "
            "e de queda, respectivamente — valores superiores indicam melhor compensação pelo risco "
            "assumido. O *Drawdown Máximo* revela a maior perda acumulada de pico a vale no período."
        ),
        "Análise Comparativa": (
            "A carteira sugerida apresenta retorno anualizado de *{ret_opt}* frente a {ret_atual} da "
            "carteira atual, *um ganho de {ganho_ret} pontos percentuais* com volatilidade de {vol_opt} "
            "(vs. {vol_atual}), demonstrando que o incremento de rentabilidade não implica risco "
            "desproporcional. O índice de Sharpe avança de {sharpe_atual} para *{sharpe_opt}*, "
            "evidenciando uma relação risco-retorno consideravelmente mais eficiente."
        ),
    },
    "volatilidade": {
        "Explicação": (
            "A volatilidade móvel de 30 dias mensura a oscilação dos retornos em janelas rolantes de "
            "um mês, expressando o risco de forma dinâmica ao longo do tempo. Períodos de maior "
            "volatilidade sinalizam incerteza elevada no mercado ou eventos específicos que afetam "
            "os ativos da carteira, enquanto ciclos de menor volatilidade indicam estabilidade relativa. "
            "A análise temporal desta métrica permite identificar regimes de risco distintos e avaliar "
            "se a carteira amplifica ou atenua a volatilidade do mercado em episódios de estresse."
        ),
        "Análise do Gráfico": (
            "O gráfico demonstra que a carteira otimizada (linha azul escuro) mantém uma "
            "*volatilidade significativamente mais suave e controlada* ao longo de todo o período analisado em comparação "
            "com a carteira atual (linha azul claro). Ambas as carteiras experimentam picos em momentos "
            "de turbulência, porém a carteira otimizada absorve esses choques com menor intensidade. "
            "*A diversificação proposta reduz significativamente a exposição a movimentos erráticos de "
            "mercado*, resultando em uma jornada de investimento mais previsível e confortável para o "
            "investidor."
        ),
    },
    "drawdown": {
        "Explicação": (
            "O gráfico de drawdown representa a queda acumulada de cada carteira em relação ao seu pico "
            "histórico mais recente. Quanto maior a profundidade e a duração do drawdown, maior o impacto "
            "sobre o patrimônio e mais exigente o processo de recuperação. Drawdowns prolongados ou severos "
            "indicam exposição elevada a fatores de risco sistemático ou concentração em ativos "
            "correlacionados. A análise comparativa permite avaliar quais alocações ofereceram maior "
            "proteção nos períodos de turbulência do mercado."
        ),
        "Análise do Gráfico": (
            "O gráfico revela uma diferença significativa entre as duas estratégias, especialmente em "
            "períodos de crise. A carteira atual (linha azul claro) experimenta *drawdowns mais profundos*, "
            "refletindo maior sensibilidade aos choques de mercado, enquanto a carteira otimizada (linha "
            "azul escuro) *limita suas perdas máximas de forma mais eficiente*. A carteira otimizada "
            "demonstra recuperação mais rápida após os períodos adversos, mantendo-se consistentemente "
            "acima da carteira atual durante momentos de estresse, evidenciando que a diversificação "
            "não apenas reduz perdas, mas também *acelera a recuperação do patrimônio*."
        ),
    },
    "correlacao": {
        "Explicação": (
            "A matriz de correlação quantifica o grau de co-movimento entre os retornos dos ativos. "
            "Células em amarelo indicam *alta correlação (≥ 0,7)*, sugerindo diversificação limitada "
            "entre os pares. Células em azul refletem *correlação moderada (0,4–0,7)*, e células em "
            "vermelho (≤ −0,3) sinalizam *ativos com tendência de movimentação oposta*, que contribuem "
            "mais efetivamente para a redução do risco por diversificação. A diagonal em azul escuro "
            "representa a autocorrelação de cada ativo consigo mesmo (sempre igual a 1). Carteiras "
            "com baixa correlação média tendem a apresentar *menor volatilidade e drawdowns mais controlados*."
        ),
    },
}


def _preencher_analise_fronteira(texto, m):
    """Fill {placeholders} in the frontier analysis text with Brazilian-formatted metrics."""
    def _pct(v):
        return f"{v * 100:.1f}".replace(".", ",") + "%"
    vals = {
        "atual_ret":         _pct(m["atual_ret"]),
        "atual_risco":       _pct(m["atual_risco"]),
        "mesmo_risco_ret":   _pct(m["mesmo_risco_ret"]),
        "mesmo_risco_risco": _pct(m["mesmo_risco_risco"]),
        "mesmo_ret_ret":     _pct(m.get("mesmo_ret_ret", m["atual_ret"])),
        "mesmo_ret_risco":   _pct(m["mesmo_ret_risco"]),
        "ganho_ret":         f"{(m['mesmo_risco_ret'] - m['atual_ret']) * 100:.2f}".replace(".", ","),
    }
    try:
        return texto.format(**vals)
    except KeyError:
        return texto


def _preencher_analise_metricas(texto, m):
    """Fill {placeholders} in the metrics analysis text with Brazilian-formatted values."""
    def _pct(v):
        return f"{v * 100:.2f}".replace(".", ",") + "%"
    def _f2(v):
        return f"{v:.2f}".replace(".", ",")
    vals = {
        "ret_atual":    _pct(m["ret_atual"]),
        "ret_opt":      _pct(m["ret_opt"]),
        "vol_atual":    _pct(m["vol_atual"]),
        "vol_opt":      _pct(m["vol_opt"]),
        "sharpe_atual": _f2(m["sharpe_atual"]) if not pd.isna(m["sharpe_atual"]) else "N/D",
        "sharpe_opt":   _f2(m["sharpe_opt"])   if not pd.isna(m["sharpe_opt"])   else "N/D",
        "dd_atual":     _pct(abs(m["dd_atual"])),
        "dd_opt":       _pct(abs(m["dd_opt"])),
        "ganho_ret":    f"{(m['ret_opt'] - m['ret_atual']) * 100:.2f}".replace(".", ","),
    }
    try:
        return texto.format(**vals)
    except KeyError:
        return texto


def _yc_layout(title="", xaxis_title="", yaxis_title="", **extra):
    base = dict(
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        font=dict(family="Raleway, Arial, sans-serif", size=11, color="#001441"),
        title=dict(text=title, font=dict(size=13, color=_YC_NAVY, family="Raleway, Arial, sans-serif")),
        xaxis=dict(
            title=xaxis_title,
            showgrid=False,
            linecolor="#9CA3AF", tickfont=dict(size=10),
        ),
        yaxis=dict(
            title=yaxis_title,
            showgrid=True, gridcolor="#EAECEF", gridwidth=1,
            linecolor="#9CA3AF", tickfont=dict(size=10),
        ),
        legend=dict(
            orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1,
            bgcolor="rgba(255,255,255,0.85)", bordercolor="#D1D5DB", borderwidth=1,
            font=dict(size=10, family="Raleway, Arial, sans-serif"),
        ),
        margin=dict(l=55, r=20, t=50, b=45),
    )
    base.update(extra)
    return base


def _make_cumret_fig(series_dict):
    aligned = _align_series_for_report(series_dict)
    fig = go.Figure()
    for i, (lbl, s) in enumerate(aligned.items()):
        fig.add_trace(go.Scatter(
            x=s.index, y=(1 + s).cumprod(), mode="lines", name=lbl,
            line=dict(color=_REPORT_PALETTE[i % len(_REPORT_PALETTE)], width=2),
        ))
    fig.update_layout(**_yc_layout(
        xaxis_title="Data", yaxis_title="Valor da Carteira",
    ))
    return fig


def _make_vol_fig(series_dict):
    aligned = _align_series_for_report(series_dict)
    fig = go.Figure()
    for i, (lbl, s) in enumerate(aligned.items()):
        fig.add_trace(go.Scatter(
            x=s.index, y=s.rolling(30).std() * np.sqrt(252), mode="lines", name=lbl,
            line=dict(color=_REPORT_PALETTE[i % len(_REPORT_PALETTE)], width=2),
        ))
    fig.update_layout(**_yc_layout(
        xaxis_title="Data", yaxis_title="Volatilidade Anualizada",
    ))
    return fig


def _make_dd_fig(series_dict):
    aligned = _align_series_for_report(series_dict)
    fig = go.Figure()
    for i, (lbl, s) in enumerate(aligned.items()):
        cum = (1 + s).cumprod()
        color = _REPORT_PALETTE[i % len(_REPORT_PALETTE)]
        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
        fig.add_trace(go.Scatter(
            x=cum.index, y=(cum / cum.cummax()) - 1, mode="lines", name=lbl, fill="tozeroy",
            line=dict(color=color, width=2),
            fillcolor=f"rgba({r},{g},{b},0.15)",
        ))
    fig.update_layout(**_yc_layout(
        xaxis_title="Data",
        yaxis=dict(
            title="Drawdown", tickformat=".1%",
            showgrid=True, gridcolor="#EAECEF", gridwidth=1,
            linecolor="#9CA3AF", tickfont=dict(size=10),
        ),
    ))
    return fig


def _make_corr_fig(returns_df):
    if len(returns_df.columns) < 2:
        return go.Figure()
    corr = returns_df.corr().round(2)
    tickers = corr.columns.tolist()
    n = len(tickers)

    def _cell_bg(row_t, col_t):
        if row_t == col_t:
            return _YC_NAVY
        v = corr.loc[row_t, col_t]
        if v >= 0.7:  return "#FEF3C7"
        if v >= 0.4:  return "#DBEAFE"
        if v <= -0.3: return "#FEE2E2"
        return "#F9FAFB"

    cell_values = [tickers]
    fill_colors = [["#EFF6FF"] * n]
    font_colors = [["#001441"] * n]

    for col_t in tickers:
        cell_values.append([f"{corr.loc[row_t, col_t]:.2f}" for row_t in tickers])
        fill_colors.append([_cell_bg(row_t, col_t) for row_t in tickers])
        font_colors.append(["#FFFFFF" if row_t == col_t else "#001441" for row_t in tickers])

    fig = go.Figure(go.Table(
        header=dict(
            values=[""] + tickers,
            fill_color=_YC_NAVY,
            font=dict(color="white", size=10, family="Raleway, Arial, sans-serif"),
            align="center", height=30,
            line=dict(color=_YC_NAVY, width=1),
        ),
        cells=dict(
            values=cell_values,
            fill_color=fill_colors,
            font=dict(color=font_colors, size=10, family="Raleway, Arial, sans-serif"),
            align="center",
            height=26,
            line=dict(color="#E5E7EB", width=0.5),
        ),
    ))
    fig.update_layout(paper_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=10, b=0))
    return fig


def _make_metrics_fig(metrics_list):
    METRICS = [
        ("Retorno Anual", "retorno_anual", ".2%"),
        ("Volatilidade", "volatilidade", ".2%"),
        ("Índice de Sharpe", "sharpe", ".2f"),
        ("Índice de Sortino", "sortino", ".2f"),
        ("Drawdown Máximo", "max_drawdown", ".2%"),
    ]
    labels = [m["label"] for m in metrics_list]
    headers = ["Métrica"] + labels
    rows = []
    for nome, key, fmt in METRICS:
        row = [nome]
        for m in metrics_list:
            v = m.get(key, np.nan)
            row.append(f"{v:{fmt}}" if not pd.isna(v) else "N/D")
        rows.append(row)
    cell_values = [list(col) for col in zip(*rows)]
    n_rows = len(METRICS)
    row_fills = ["#FFFFFF" if i % 2 == 0 else "#F3F4F6" for i in range(n_rows)]
    fill_colors = [row_fills] * len(headers)
    fig = go.Figure(go.Table(
        header=dict(
            values=headers, fill_color=_YC_NAVY,
            font=dict(color="white", size=12, family="Raleway, Arial, sans-serif"),
            align=["left"] + ["center"] * len(labels), height=38,
            line=dict(color=_YC_NAVY, width=0),
        ),
        cells=dict(
            values=cell_values, fill_color=fill_colors,
            align=["left"] + ["center"] * len(labels),
            font=dict(size=12, family="Raleway, Arial, sans-serif", color="#001441"),
            height=34, line=dict(color="#E5E7EB", width=0.5),
        ),
    ))
    fig.update_layout(paper_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=0, b=0))
    return fig


def _make_allocation_table_fig(weights_atual, weights_sug):
    keys_a = list(weights_atual.keys()) if weights_atual else []
    keys_s = list(weights_sug.keys()) if weights_sug else []
    all_tickers = sorted(set(keys_a + keys_s))
    if not all_tickers:
        return None
    has_sug = bool(weights_sug)
    tickers_col = [symbol_info.get(t, t) for t in all_tickers]
    atual_col = [f"{weights_atual.get(t, 0):.2%}" for t in all_tickers]
    if has_sug:
        sug_col = [f"{weights_sug.get(t, 0):.2%}" for t in all_tickers]
        diff_col = [f"{weights_sug.get(t, 0) - weights_atual.get(t, 0):+.2%}" for t in all_tickers]
        headers = ["Ativo", "Atual %", "Sugerido %", "Diferença %"]
        cell_values = [tickers_col, atual_col, sug_col, diff_col]
        col_align = ["left", "right", "right", "right"]
    else:
        headers = ["Ativo", "Atual %"]
        cell_values = [tickers_col, atual_col]
        col_align = ["left", "right"]
    n_rows = len(all_tickers)
    row_fills = ["#FFFFFF" if i % 2 == 0 else "#F3F4F6" for i in range(n_rows)]
    fill_colors = [row_fills] * len(headers)
    fig = go.Figure(go.Table(
        header=dict(
            values=headers, fill_color=_YC_NAVY,
            font=dict(color="white", size=12, family="Raleway, Arial, sans-serif"),
            align=["left"] + ["right"] * (len(headers) - 1), height=38,
            line=dict(color=_YC_NAVY, width=0),
        ),
        cells=dict(
            values=cell_values, fill_color=fill_colors, align=col_align,
            font=dict(size=12, family="Raleway, Arial, sans-serif", color="#001441"),
            height=34, line=dict(color="#E5E7EB", width=0.5),
        ),
    ))
    fig.update_layout(paper_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=0, b=0))
    return fig


def _gerar_relatorio_pptx(client_name, portfolio_value, series_dict, returns_df, weights_atual, weights_sug, benchmark_label, frontier_fig=None, slide_texts=None, template_path="analise.pptx"):
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(template_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    # Proportional scaling helpers — calibrated for A4 portrait (8.264 × 11.694 in).
    # All Inches/Pt values below are in A4 equivalents; _sw/_sh/_sf scale them to any
    # slide size so the layout stays correct even when the template changes dimensions.
    _REF_W, _REF_H = 7556500.0, 10693400.0
    _sw = lambda a4in: int(a4in * 914400 * slide_w / _REF_W)
    _sh = lambda a4in: int(a4in * 914400 * slide_h / _REF_H)
    _sf = lambda a4pt: Pt(max(6, round(a4pt * slide_h / _REF_H)))

    margin_l  = _sw(0.65)
    content_w = slide_w - 2 * margin_l
    blank_layout = prs.slide_layouts[6]
    n_template_slides = len(prs.slides)

    title_top = _sh(0.85)
    title_h   = _sh(0.7)
    chart_top = title_top + title_h + _sh(0.15)
    chart_h   = int(slide_h / 3)

    # Fixed image resolution — kaleido renders in pixels, independent of slide size
    px_w, px_h = 900, 400

    # --- Capa: insere nome do cliente (canto inferior direito) ---
    capa = prs.slides[0]
    name_w = _sw(5.16)
    name_h = _sh(0.65)
    txb = capa.shapes.add_textbox(slide_w - margin_l - name_w, _sh(8.4), name_w, name_h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = client_name
    run.font.name = "Raleway"
    run.font.size = _sf(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.RIGHT

    # Padroniza label da carteira otimizada antes de construir métricas e figuras
    series_dict = {
        ("Carteira Otimizada" if k.startswith("Otimizado") else k): v
        for k, v in series_dict.items()
    }

    # --- Métricas (exclui benchmark) ---
    rf_metric = 0.0
    metrics_list = [
        _calcular_metricas(s, rf_metric, lbl)
        for lbl, s in series_dict.items()
        if lbl != benchmark_label
    ]
    if not metrics_list:
        metrics_list = [_calcular_metricas(s, rf_metric, lbl) for lbl, s in series_dict.items()]

    # --- Figuras (apenas gráficos; tabelas são geradas como objetos PPTX nativos) ---
    fig_cum = _make_cumret_fig(series_dict)
    fig_vol = _make_vol_fig(series_dict)
    fig_dd  = _make_dd_fig(series_dict)

    def _add_title(slide, text):
        txb2 = slide.shapes.add_textbox(margin_l, title_top, content_w, title_h)
        tf2 = txb2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = text
        run2.font.size = _sf(22)
        run2.font.bold = True
        run2.font.name = "Raleway"
        run2.font.color.rgb = RGBColor(0x00, 0x14, 0x41)

    def _add_image(slide, fig, left, top, width, height, px_w=900, px_h=400):
        png = _fig_to_png_bytes(fig, width=px_w, height=px_h)
        slide.shapes.add_picture(io.BytesIO(png), left, top, width, height)

    def _pptx_table(slide, headers, rows_data, left, top, width,
                    row_h=None, header_h=None,
                    col_aligns=None, col_widths=None,
                    cell_fills=None, cell_font_colors=None):
        if row_h    is None: row_h    = _sh(0.30)
        if header_h is None: header_h = _sh(0.35)
        """Create a native PPTX table. Returns (shape, total_height_in_EMU)."""
        n_cols  = len(headers)
        n_data  = len(rows_data)
        total_h = int(header_h) + int(row_h) * n_data
        shape   = slide.shapes.add_table(n_data + 1, n_cols, left, top, width, total_h)
        tbl     = shape.table
        if col_widths:
            for ci, cw in enumerate(col_widths):
                tbl.columns[ci].width = int(cw)
        _NAV = RGBColor(0x00, 0x14, 0x41)
        _WHT = RGBColor(0xFF, 0xFF, 0xFF)
        _DRK = RGBColor(0x00, 0x14, 0x41)
        _BG1 = RGBColor(0xFF, 0xFF, 0xFF)
        _BG2 = RGBColor(0xF3, 0xF4, 0xF6)
        for ci, hdr in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _NAV
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(hdr)
            run.font.size = _sf(10)
            run.font.bold = True
            run.font.name = "Raleway"
            run.font.color.rgb = _WHT
        for ri, row in enumerate(rows_data):
            default_bg = _BG1 if ri % 2 == 0 else _BG2
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + 1, ci)
                bg = (cell_fills[ri][ci]
                      if cell_fills and ri < len(cell_fills) and ci < len(cell_fills[ri])
                         and cell_fills[ri][ci] is not None
                      else default_bg)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                p = cell.text_frame.paragraphs[0]
                if col_aligns and ci < len(col_aligns):
                    a = col_aligns[ci]
                    p.alignment = (PP_ALIGN.RIGHT  if a == "right"
                                   else PP_ALIGN.CENTER if a == "center"
                                   else PP_ALIGN.LEFT)
                else:
                    p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = str(val)
                run.font.size = _sf(10)
                run.font.name = "Raleway"
                run.font.color.rgb = (
                    cell_font_colors[ri][ci]
                    if cell_font_colors and ri < len(cell_font_colors)
                       and ci < len(cell_font_colors[ri])
                       and cell_font_colors[ri][ci] is not None
                    else _DRK
                )
        return shape, total_h

    def _add_logo_slide(slide):
        logo_path = "icone_final.png"
        if not os.path.exists(logo_path):
            return
        try:
            with Image.open(logo_path) as _img:
                _lw_px, _lh_px = _img.size
            _prop = _lw_px / _lh_px if _lh_px > 0 else 1
            _lw = _sh(0.55)
            _lh = int(_lw / _prop)
            if _lh > _sh(0.55):
                _lh = _sh(0.55)
                _lw = int(_lh * _prop)
            _lx = slide_w - _sw(0.80) - _lw
            _ly = _sh(10.55)
            slide.shapes.add_picture(logo_path, _lx, _ly, width=_lw, height=_lh)
        except Exception:
            pass

    def _add_texto_slide(slide, texto, top):
        """Add a justified body text box below a chart. Supports *text* for bold."""
        if not texto or not texto.strip():
            return
        paras = [p.strip() for p in texto.split("\n") if p.strip()]
        n_linhas = sum(max(1, math.ceil(len(p) / 88)) for p in paras)
        altura = min(
            int(_sh(0.24) * n_linhas + _sh(0.14) * len(paras)),
            int(slide_h - top - _sh(0.15)),
        )
        if altura < int(_sh(0.18)):
            return
        txb = slide.shapes.add_textbox(margin_l, top, content_w, altura)
        tf  = txb.text_frame
        tf.clear()
        tf.word_wrap    = True
        tf.margin_left  = _sw(0.02)
        tf.margin_right = _sw(0.02)
        tf.margin_top   = _sh(0.02)
        tf.margin_bottom = _sh(0.02)
        _TCOL = RGBColor(0x00, 0x14, 0x41)
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment    = PP_ALIGN.JUSTIFY
            p.line_spacing = _sf(16)
            p.space_after  = _sf(10)
            for j, seg in enumerate(re.split(r'\*([^*]+)\*', para)):
                if not seg:
                    continue
                run = p.add_run()
                run.text           = seg
                run.font.name      = "Raleway"
                run.font.size      = _sf(12)
                run.font.bold      = (j % 2 == 1)
                run.font.color.rgb = _TCOL

    _st = slide_texts or {}

    # --- Slide Fronteira Eficiente (apenas Tab 1) ---
    if frontier_fig is not None:
        sl_fe = prs.slides.add_slide(blank_layout)
        _add_logo_slide(sl_fe)
        _add_title(sl_fe, "Fronteira Eficiente")
        _add_image(sl_fe, frontier_fig, margin_l, chart_top, content_w, chart_h,
                   px_w=px_w, px_h=px_h)
        _add_texto_slide(sl_fe, _st.get("fronteira", ""), chart_top + chart_h + _sh(0.15))

    # --- Slide A: Métricas de Desempenho + Retorno Acumulado ---
    sl_a = prs.slides.add_slide(blank_layout)
    _add_logo_slide(sl_a)
    _add_title(sl_a, "Métricas de Desempenho")
    _MET_DEFS = [
        ("Retorno Anual",     "retorno_anual", ".2%"),
        ("Volatilidade",      "volatilidade",  ".2%"),
        ("Índice de Sharpe",  "sharpe",        ".2f"),
        ("Índice de Sortino", "sortino",        ".2f"),
        ("Drawdown Máximo",   "max_drawdown",  ".2%"),
    ]
    met_headers = ["Métrica"] + [m["label"] for m in metrics_list]
    met_rows = []
    for nome, key, fmt in _MET_DEFS:
        row = [nome]
        for m in metrics_list:
            v = m.get(key, float("nan"))
            row.append(f"{v:{fmt}}" if not pd.isna(v) else "N/D")
        met_rows.append(row)
    n_met_cols = len(met_headers)
    met_cw = ([int(content_w * 0.40)]
              + [int(content_w * 0.60 / max(1, n_met_cols - 1))] * (n_met_cols - 1))
    _, met_tbl_h = _pptx_table(
        sl_a, met_headers, met_rows, margin_l, chart_top, content_w,
        row_h=_sh(0.30), header_h=_sh(0.35),
        col_aligns=["left"] + ["center"] * (n_met_cols - 1),
        col_widths=met_cw,
    )
    cum_top = chart_top + met_tbl_h + _sh(0.2)
    cum_h   = int(slide_h * 0.28)
    _add_image(sl_a, fig_cum, margin_l, cum_top, content_w, cum_h,
               px_w=px_w, px_h=max(240, px_h))
    _add_texto_slide(sl_a, _st.get("metricas", ""), cum_top + cum_h + _sh(0.15))

    # --- Slide B: Volatilidade Móvel ---
    sl_b = prs.slides.add_slide(blank_layout)
    _add_logo_slide(sl_b)
    _add_title(sl_b, "Volatilidade Móvel (30 Dias)")
    _add_image(sl_b, fig_vol, margin_l, chart_top, content_w, chart_h,
               px_w=px_w, px_h=px_h)
    _add_texto_slide(sl_b, _st.get("volatilidade", ""), chart_top + chart_h + _sh(0.15))

    # --- Slide C: Drawdown ---
    sl_c = prs.slides.add_slide(blank_layout)
    _add_logo_slide(sl_c)
    _add_title(sl_c, "Drawdown da Carteira")
    _add_image(sl_c, fig_dd, margin_l, chart_top, content_w, chart_h,
               px_w=px_w, px_h=px_h)
    _add_texto_slide(sl_c, _st.get("drawdown", ""), chart_top + chart_h + _sh(0.15))

    # --- Slide D: Matriz de Correlação ---
    sl_d = prs.slides.add_slide(blank_layout)
    _add_logo_slide(sl_d)
    _add_title(sl_d, "Matriz de Correlação dos Ativos")
    _corr   = returns_df.corr().round(2)
    _tkrs_c = _corr.columns.tolist()
    corr_headers = [" "] + [symbol_info.get(t, t) for t in _tkrs_c]
    corr_rows, corr_fills, corr_fcolors = [], [], []
    for rt in _tkrs_c:
        row_vals  = [symbol_info.get(rt, rt)]
        row_fills = [RGBColor(0xEF, 0xF6, 0xFF)]
        row_fc    = [None]
        for ct in _tkrs_c:
            v = _corr.loc[rt, ct]
            row_vals.append(f"{v:.2f}")
            if rt == ct:
                row_fills.append(RGBColor(0x00, 0x14, 0x41))
                row_fc.append(RGBColor(0xFF, 0xFF, 0xFF))
            elif v >= 0.7:
                row_fills.append(RGBColor(0xFE, 0xF3, 0xC7)); row_fc.append(None)
            elif v >= 0.4:
                row_fills.append(RGBColor(0xDB, 0xEA, 0xFE)); row_fc.append(None)
            elif v <= -0.3:
                row_fills.append(RGBColor(0xFE, 0xE2, 0xE2)); row_fc.append(None)
            else:
                row_fills.append(None); row_fc.append(None)
        corr_rows.append(row_vals)
        corr_fills.append(row_fills)
        corr_fcolors.append(row_fc)
    n_tkrs_c = len(_tkrs_c)
    corr_cw  = ([int(content_w * 0.18)]
                + [int(content_w * 0.82 / max(1, n_tkrs_c))] * n_tkrs_c)
    corr_shape, _ = _pptx_table(
        sl_d, corr_headers, corr_rows, margin_l, chart_top, content_w,
        row_h=_sh(0.28), header_h=_sh(0.32),
        col_aligns=["left"] + ["center"] * n_tkrs_c,
        col_widths=corr_cw,
        cell_fills=corr_fills,
        cell_font_colors=corr_fcolors,
    )
    _add_texto_slide(sl_d, _st.get("correlacao", ""), corr_shape.top + corr_shape.height + _sh(0.35))

    # --- Slide E: Alocação (opcional) ---
    _tkrs_a = sorted(set(
        (list(weights_atual.keys()) if weights_atual else []) +
        (list(weights_sug.keys())   if weights_sug   else [])
    ))
    if _tkrs_a:
        sl_e = prs.slides.add_slide(blank_layout)
        _add_logo_slide(sl_e)
        _add_title(sl_e, "Movimentação da Carteira")
        _GRN = RGBColor(0x05, 0x96, 0x69)
        _RED = RGBColor(0xDC, 0x26, 0x26)
        if weights_sug:
            alloc_hdrs = ["Ativo", "Atual %", "Sugerido %", "Diferença %", "Movimentação (R$)"]
            alloc_rows, alloc_fc = [], []
            for t in _tkrs_a:
                wa  = weights_atual.get(t, 0)
                ws  = weights_sug.get(t, 0)
                dif = ws - wa
                fin = dif * portfolio_value
                alloc_rows.append([
                    symbol_info.get(t, t),
                    f"{wa:.2%}", f"{ws:.2%}",
                    f"{dif:+.2%}",
                    _brl(fin),
                ])
                dfc = _GRN if dif > 0.0005 else (_RED if dif < -0.0005 else None)
                ffc = _GRN if fin > 0      else (_RED if fin < 0      else None)
                alloc_fc.append([None, None, None, dfc, ffc])
            alloc_cw = [int(content_w * p) for p in [0.30, 0.14, 0.14, 0.15, 0.27]]
            _pptx_table(
                sl_e, alloc_hdrs, alloc_rows, margin_l, chart_top, content_w,
                row_h=_sh(0.30), header_h=_sh(0.35),
                col_aligns=["left", "right", "right", "right", "right"],
                col_widths=alloc_cw,
                cell_font_colors=alloc_fc,
            )
        else:
            alloc_hdrs = ["Ativo", "Atual %"]
            alloc_rows = [[symbol_info.get(t, t), f"{weights_atual.get(t, 0):.2%}"]
                          for t in _tkrs_a]
            _pptx_table(sl_e, alloc_hdrs, alloc_rows, margin_l, chart_top, content_w)

    # --- Move último slide do template (contato) para o final ---
    xml_list = prs.slides._sldIdLst
    contact_elem = xml_list[n_template_slides - 1]
    xml_list.remove(contact_elem)
    xml_list.append(contact_elem)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _render_report_ui(tab_key):
    """Renders the PPTX report generation UI for either tab."""
    series_key   = f"{tab_key}_report_series"
    returns_key  = f"{tab_key}_report_returns"
    wa_key       = f"{tab_key}_weights_atual"
    ws_key       = f"{tab_key}_weights_sug"
    bm_key       = f"{tab_key}_benchmark"
    buf_key      = f"{tab_key}_report_buf"
    name_key     = f"{tab_key}_report_name"

    if not st.session_state.get(series_key):
        return

    _MODELOS_RELATORIO = {"Relatório 1": "analise.pptx", "Relatório 2": "analise1.pptx"}

    st.markdown("---")
    with st.expander("Gerar Relatório PPTX", expanded=False):
        r_modelo = st.radio(
            "Modelo do relatório:",
            list(_MODELOS_RELATORIO.keys()),
            key=f"{tab_key}_r_modelo",
            horizontal=True,
        )
        r_client = st.text_input("Nome do cliente", key=f"{tab_key}_r_client")
        r_value  = st.number_input(
            "Valor total da carteira (R$)",
            min_value=0.0, value=1_000_000.0, step=10_000.0,
            key=f"{tab_key}_r_value",
        )

        st.markdown("**Textos por slide** *(opcional — aparecem abaixo de cada gráfico/tabela)*")

        slide_texts = {}

        # --- Sections with morningcall-style multi-paragraph UI ---
        _SECS_MULTI = [
            ("fronteira",    "Fronteira Eficiente",          tab_key == "t1"),
            ("metricas",     "Métricas de Desempenho",       True),
            ("volatilidade", "Volatilidade Móvel (30 Dias)", True),
            ("drawdown",     "Drawdown da Carteira",         True),
        ]
        for _sk, _slabel, _show in _SECS_MULTI:
            if not _show:
                continue
            st.markdown(f"**{_slabel}**")
            _n = int(st.number_input(
                "Parágrafos", min_value=1, max_value=5, value=2, step=1,
                key=f"{tab_key}_np_{_sk}",
            ))
            _opts = ["Sem texto"] + list(_TEXTOS_PADRAO[_sk].keys())
            _paras = []
            for _pi in range(_n):
                _sel = st.selectbox(
                    f"Parágrafo {_pi + 1}",
                    _opts,
                    key=f"{tab_key}_sel_{_sk}_{_pi}",
                )
                if _sel != "Sem texto":
                    _tb = _TEXTOS_PADRAO[_sk][_sel]
                    if "{atual_ret}" in _tb:
                        _m = st.session_state.get(f"{tab_key}_frontier_metrics", {})
                        if _m:
                            _tb = _preencher_analise_fronteira(_tb, _m)
                    elif "{ret_atual}" in _tb:
                        _m = st.session_state.get(f"{tab_key}_metrics_data", {})
                        if _m:
                            _tb = _preencher_analise_metricas(_tb, _m)
                    _txt = st.text_area(
                        f"Texto — parágrafo {_pi + 1}",
                        value=_tb,
                        key=f"{tab_key}_txt_{_sk}_{_pi}_{_sel}",
                        height=110,
                        label_visibility="collapsed",
                    )
                else:
                    _txt = st.text_area(
                        f"Texto — parágrafo {_pi + 1}",
                        value="",
                        key=f"{tab_key}_blank_{_sk}_{_pi}",
                        height=90,
                        label_visibility="collapsed",
                    )
                _paras.append(_txt)
            slide_texts[_sk] = "\n".join(p for p in _paras if p.strip())
            st.markdown("---")

        # --- Correlação: single text area ---
        st.markdown("**Matriz de Correlação**")
        _opts_corr = ["Sem texto"] + list(_TEXTOS_PADRAO["correlacao"].keys())
        _sel_corr = st.selectbox(
            "Matriz de Correlação",
            _opts_corr,
            key=f"{tab_key}_sel_correlacao",
            label_visibility="collapsed",
        )
        if _sel_corr != "Sem texto":
            slide_texts["correlacao"] = st.text_area(
                "Texto editável",
                value=_TEXTOS_PADRAO["correlacao"][_sel_corr],
                key=f"{tab_key}_txt_correlacao_{_sel_corr}",
                height=110,
                label_visibility="collapsed",
            )
        else:
            slide_texts["correlacao"] = ""

        if st.button("Gerar Relatório", key=f"{tab_key}_r_btn"):
            if r_client.strip():
                with st.spinner("Gerando relatório PPTX…"):
                    report_bytes = _gerar_relatorio_pptx(
                        r_client.strip(),
                        r_value,
                        st.session_state[series_key],
                        st.session_state[returns_key],
                        st.session_state.get(wa_key) or {},
                        st.session_state.get(ws_key),
                        st.session_state.get(bm_key),
                        frontier_fig=st.session_state.get(f"{tab_key}_frontier_fig"),
                        slide_texts=slide_texts,
                        template_path=_MODELOS_RELATORIO[r_modelo],
                    )
                st.session_state[buf_key]  = report_bytes
                st.session_state[name_key] = r_client.strip()
            else:
                st.error("Digite o nome do cliente.")

        if st.session_state.get(buf_key):
            safe = st.session_state[name_key].replace(" ", "_")
            st.download_button(
                "Baixar Relatório PPTX",
                data=st.session_state[buf_key],
                file_name=f"relatorio_{safe}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key=f"{tab_key}_r_download",
            )


def show_comparador():
    try:
        logo = Image.open("logo_final.png")
        st.image(logo, use_container_width=False, width=800)

        if "synthetic_assets" not in st.session_state:
            st.session_state["synthetic_assets"] = {}

        tab1, tab2 = st.tabs(["Comparador com Fronteira Eficiente", "Comparação Simples"])

        ticker_options = [f"{desc} [{sym}]" for sym, desc in symbol_info.items()]
        ticker_lookup = {f"{desc} [{sym}]": sym for sym, desc in symbol_info.items()}
        # Ativos sintéticos entram como opções do editor — o peso conta para os 100%
        for _sn, _si in st.session_state["synthetic_assets"].items():
            _sd = f"{_sn} (Sintético β={_si['beta']:.2g}×{_si['ref_ticker']}) [{_sn}]"
            ticker_options.append(_sd)
            ticker_lookup[_sd] = _sn
        default_symbols = ["IFIX", "IMA-B 5", "IDA-IPCA", "GLD", "USDBRL=X", "DIVO11.SA", "BTC-USD", "SPY", "IHFA", "CDI"]
        default_weights = [0.10] * 10

        # =============================================================
        # TAB 1 — Comparador com Fronteira Eficiente
        # =============================================================
        with tab1:
            st.subheader("Compare sua Carteira com uma Otimizada")

            if "comparison_ready" not in st.session_state:
                st.session_state["comparison_ready"] = False
            if "comparison_running" not in st.session_state:
                st.session_state["comparison_running"] = False

            if st.session_state["comparison_ready"]:
                col_back, _ = st.columns([1, 5])
                with col_back:
                    if st.button("← Voltar à configuração", key="back_to_setup_comparador"):
                        st.session_state["comparison_ready"] = False
                        st.session_state["comparison_running"] = False
                        st.rerun()

            if not st.session_state["comparison_ready"]:
                st.markdown("Selecione os ativos da sua carteira atual e seus respectivos pesos.")
                _render_synth_expander("t1", ticker_options, ticker_lookup)

                _currency_t1 = st.session_state.get("currency_comparador", "Real")
                user_portfolio, weight_sum = _portfolio_editor(
                    default_symbols, default_weights, ticker_options, ticker_lookup, "t1_atual",
                    currency_choice=_currency_t1,
                )

                _peso_invalido = abs(weight_sum - 100) > 0.01
                if _peso_invalido:
                    st.error("A alocação total deve somar 100%. Por favor, ajuste os pesos.")
                if user_portfolio["Weight"].sum() > 0:
                    user_portfolio["Weight"] /= user_portfolio["Weight"].sum()

                available_symbols = sorted(symbol_info.keys(), key=lambda x: symbol_info.get(x, x))
                _synth_keys = list(st.session_state.get("synthetic_assets", {}).keys())
                available_symbols_with_synth = available_symbols + _synth_keys
                current_symbols = user_portfolio["Ticker"].dropna().tolist()
                extra_default = [s for s in ["SPY", "BND", "GLD"] if s not in current_symbols and s in symbol_info]
                # Restore previously selected extra assets (including synthetics) after "Voltar"
                _saved_extra = st.session_state.get("extra_assets")
                _extra_init = (
                    [s for s in _saved_extra if s in available_symbols_with_synth]
                    if _saved_extra is not None else extra_default
                )
                _synth_info_map = st.session_state.get("synthetic_assets", {})
                extra_assets = st.multiselect(
                    "Selecione ativos candidatos adicionais:",
                    options=available_symbols_with_synth,
                    default=_extra_init,
                    format_func=lambda x: (
                        f"{x} (Sintético β={_synth_info_map[x]['beta']:.2g}×{_synth_info_map[x]['ref_ticker']}) [{x}]"
                        if x in _synth_info_map
                        else f"{symbol_info.get(x, x)} [{x}]"
                    ),
                    key="extra_assets_comparador",
                )

                candidate_tickers = list(dict.fromkeys(user_portfolio["Ticker"].dropna().unique().tolist() + extra_assets))

                _currencies = ["Real", "Dollar", "Euro"]
                _saved_cur = st.session_state.get("currency_choice", "Real")
                _cur_idx = _currencies.index(_saved_cur) if _saved_cur in _currencies else 0
                currency_choice = st.selectbox("Selecione a Moeda", _currencies, index=_cur_idx, key="currency_comparador")

                _risk_opts = ["Mean-Variance", "Mean-Semivariance", "Mean Absolute Deviation (MAD)", "CVaR / Expected Shortfall"]
                _saved_risk = st.session_state.get("risk_measure_comparador", "Mean-Variance")
                _risk_idx = _risk_opts.index(_saved_risk) if _saved_risk in _risk_opts else 0
                risk_measure_comparador = st.selectbox(
                    "Método de Construção da Fronteira",
                    _risk_opts,
                    index=_risk_idx,
                    key="risk_measure_comparador_widget",
                )

                confidence_level_comparador = 0.95
                if risk_measure_comparador == "CVaR / Expected Shortfall":
                    _saved_conf = st.session_state.get("confidence_level_comparador", 0.95)
                    confidence_level_comparador = st.slider(
                        "Nível de Confiança", min_value=0.90, max_value=0.99,
                        value=float(_saved_conf), step=0.01,
                        key="confidence_level_comparador_widget",
                    )

                benchmark_options = ["Nenhum"] + list(benchmarks.keys())
                _saved_bench = st.session_state.get("benchmark_choice_comparador", "Nenhum")
                _bench_idx = benchmark_options.index(_saved_bench) if _saved_bench in benchmark_options else 0
                benchmark_choice_comparador = st.selectbox(
                    "Selecione o Benchmark para Comparação:", benchmark_options,
                    index=_bench_idx, key="benchmark_comparador"
                )

                g_min_t1, g_max_t1, per_asset_t1, groups_t1 = _render_constraints_expander(
                    candidate_tickers, "comparador_t1"
                )

                today = pd.Timestamp.today().date()
                # Pré-carrega índices BR se algum foi selecionado
                br_setup = [s for s in candidate_tickers if s in symbol_br_indices]
                br_idx_setup = _get_br_slice(start="2000-01-01", end=str(today)) if br_setup else None

                _synth_setup = st.session_state.get("synthetic_assets", {})
                price_data = pd.DataFrame()
                valid_symbols = []
                for symbol in candidate_tickers:
                    if symbol in symbol_br_indices:
                        if br_idx_setup is not None and symbol in br_idx_setup.columns:
                            valid_symbols.append(symbol)
                        else:
                            st.warning(f"⚠️ Dados não disponíveis para: {symbol}")
                    elif symbol in _synth_setup:
                        valid_symbols.append(symbol)
                    else:
                        close_data = safe_download(symbol, start="2000-01-01", end=today)
                        if not close_data.empty:
                            price_data[symbol] = close_data
                            valid_symbols.append(symbol)
                        else:
                            st.warning(f"⚠️ Falha ao baixar dados para: {symbol}")

                if len(valid_symbols) < 2:
                    st.error("Poucos ativos válidos para executar a análise. Por favor, ajuste sua seleção.")
                else:
                    synth_valid = [s for s in valid_symbols if s in _synth_setup]
                    combined_setup = price_data.copy() if not price_data.empty else pd.DataFrame()
                    if br_idx_setup is not None:
                        br_cols = [s for s in valid_symbols if s in symbol_br_indices and s in br_idx_setup.columns]
                        for _bc in br_cols:
                            combined_setup[_bc] = br_idx_setup[_bc]
                    if not combined_setup.empty:
                        _common = combined_setup.dropna()
                        first_date = _common.index.min() if not _common.empty else combined_setup.index.min()
                        last_date = _common.index.max() if not _common.empty else combined_setup.index.max()
                    elif synth_valid:
                        _ref = _synth_setup[synth_valid[0]]["ref_ticker"]
                        _ref_data = safe_download(_ref, start="2000-01-01", end=today)
                        if not _ref_data.empty:
                            first_date = _ref_data.index.min()
                            last_date = _ref_data.index.max()
                        else:
                            first_date = pd.Timestamp("2010-01-01")
                            last_date = pd.Timestamp(today)
                    else:
                        first_date = pd.Timestamp("2010-01-01")
                        last_date = pd.Timestamp(today)

                    _saved_start = st.session_state.get("start_date")
                    _saved_end   = st.session_state.get("end_date")
                    _start_val = (
                        _saved_start if _saved_start is not None and first_date.date() <= _saved_start <= last_date.date()
                        else first_date.date()
                    )
                    _end_val = (
                        _saved_end if _saved_end is not None and first_date.date() <= _saved_end <= last_date.date()
                        else last_date.date()
                    )
                    col1, col2 = st.columns(2)
                    with col1:
                        start_date = st.date_input("Data Inicial", value=_start_val, min_value=first_date.date(), max_value=last_date.date(), key="t1_start")
                    with col2:
                        end_date = st.date_input("Data Final", value=_end_val, min_value=first_date.date(), max_value=last_date.date(), key="t1_end")

                    if st.button("Executar Comparação", key="run_comparison_comparador", disabled=_peso_invalido):
                        st.session_state["user_portfolio"] = user_portfolio
                        st.session_state["start_date"] = start_date
                        st.session_state["end_date"] = end_date
                        st.session_state["comparison_ready"] = True
                        st.session_state["comparison_running"] = True
                        st.session_state["extra_assets"] = extra_assets
                        st.session_state["benchmark_choice_comparador"] = benchmark_choice_comparador
                        st.session_state["risk_measure_comparador"] = risk_measure_comparador
                        st.session_state["confidence_level_comparador"] = confidence_level_comparador
                        st.session_state["currency_choice"] = currency_choice
                        st.session_state["alloc_g_min"] = g_min_t1
                        st.session_state["alloc_g_max"] = g_max_t1
                        st.session_state["alloc_per_asset"] = per_asset_t1
                        st.session_state["alloc_groups"] = groups_t1
                        # Estes data_editors não são renderizados enquanto a tela de
                        # resultados estiver visível. Remover as chaves explicitamente
                        # (em vez de depender do momento em que o Streamlit faz essa
                        # limpeza sozinho) garante que, ao clicar em "Voltar", o editor
                        # sempre seja reidratado a partir do snapshot — inclusive em
                        # ciclos repetidos de Executar → Voltar.
                        for _k in ("t1_atual_editor", "t1_atual_editor_val", "alloc_override_comparador_t1", "alloc_groups_comparador_t1"):
                            st.session_state.pop(_k, None)
                        st.rerun()

            else:
                prog = st.progress(0, text="Iniciando comparação...")
                user_portfolio = st.session_state["user_portfolio"]
                extra_assets = st.session_state.get("extra_assets", [])
                start_date = st.session_state["start_date"]
                end_date = st.session_state["end_date"]
                benchmark_choice_comparador = st.session_state.get("benchmark_choice_comparador", "Nenhum")
                risk_measure_comparador = st.session_state.get("risk_measure_comparador", "Mean-Variance")
                confidence_level_comparador = st.session_state.get("confidence_level_comparador", 0.95)
                currency_choice = st.session_state.get("currency_choice", "Real")
                today = pd.Timestamp.today().date()
                prog.progress(10, text="Preparando universo de ativos...")

                portfolio_tickers = user_portfolio["Ticker"].dropna().tolist()
                tickers = list(dict.fromkeys(portfolio_tickers + extra_assets))
                weights_current = user_portfolio.set_index("Ticker")["Weight"].to_dict()
                benchmark_ticker_comparador = benchmarks.get(benchmark_choice_comparador)

                # Separa tickers yfinance dos índices BR e dos sintéticos
                _synth = st.session_state.get("synthetic_assets", {})
                yf_tickers = [t for t in tickers if t not in symbol_br_indices and t not in _synth]
                br_tickers = [t for t in tickers if t in symbol_br_indices]
                synth_tickers_t1 = [t for t in tickers if t in _synth]

                raw_data = yf.download(yf_tickers, start=start_date, end=end_date, group_by="ticker") if yf_tickers else pd.DataFrame()
                prog.progress(30, text="Baixando dados da carteira...")
                if not raw_data.empty and isinstance(raw_data.columns, pd.MultiIndex):
                    price_data = raw_data.xs("Close", axis=1, level=1)[yf_tickers].dropna()
                elif not raw_data.empty:
                    price_data = raw_data[[t for t in yf_tickers if t in raw_data.columns]].dropna()
                else:
                    price_data = pd.DataFrame()

                # Benchmark
                if benchmark_ticker_comparador is not None and benchmark_ticker_comparador not in BR_BENCH_TICKERS:
                    bm_price = safe_download(benchmark_ticker_comparador, start=start_date, end=end_date)
                    if isinstance(bm_price, pd.DataFrame):
                        bm_price = bm_price.squeeze()
                    bm_price = pd.Series(bm_price).dropna()
                    if bm_price.empty:
                        st.warning(f"⚠️ Falha ao baixar dados do benchmark: {benchmark_choice_comparador}")
                    benchmark_returns = pd.Series(bm_price.pct_change().dropna()).dropna()
                elif benchmark_ticker_comparador in BR_BENCH_TICKERS:
                    br_idx_bm = _get_br_slice(start=str(start_date), end=str(end_date))
                    if benchmark_ticker_comparador in br_idx_bm.columns:
                        benchmark_returns = br_idx_bm[benchmark_ticker_comparador].dropna()
                    elif benchmark_ticker_comparador == "CDI":
                        try:
                            benchmark_returns = _load_cdi_bcb(start=str(start_date), end=str(end_date))
                        except Exception as _cdi_e:
                            st.warning(f"⚠️ Não foi possível carregar o CDI do Banco Central: {_cdi_e}")
                            benchmark_returns = pd.Series(dtype=float)
                    else:
                        st.warning(f"⚠️ Índice '{benchmark_choice_comparador}' não disponível no período selecionado.")
                        benchmark_returns = pd.Series(dtype=float)
                else:
                    benchmark_returns = pd.Series(dtype=float)
                prog.progress(45, text="Baixando dados do benchmark...")

                returns = price_data.pct_change().dropna() if not price_data.empty else pd.DataFrame()
                prog.progress(60, text="Calculando retornos e risco...")

                benchmark_returns = _ajustar_cambio_benchmark(benchmark_returns, benchmark_ticker_comparador, currency_choice, today)
                returns = _ajustar_cambio(returns, currency_choice, today)

                # Mescla retornos dos índices BR
                if br_tickers:
                    br_idx_res = _get_br_slice(start=str(start_date), end=str(end_date))
                    if returns.empty:
                        returns = br_idx_res[[t for t in br_tickers if t in br_idx_res.columns]].dropna()
                    else:
                        for t in br_tickers:
                            if t in br_idx_res.columns:
                                returns[t] = br_idx_res[t]
                        avail_br = [t for t in br_tickers if t in returns.columns]
                        if avail_br:
                            returns = returns.dropna(subset=avail_br)

                # Adiciona retornos dos ativos sintéticos
                for st_name in synth_tickers_t1:
                    _si = _synth[st_name]
                    _sr = _compute_synthetic_returns(_si["ref_ticker"], start_date, end_date)
                    if _sr is not None:
                        returns[st_name] = _si["beta"] * _sr
                    else:
                        st.warning(f"⚠️ Não foi possível calcular retornos para o ativo sintético '{st_name}' (referência: {_si['ref_ticker']}).")

                weights_vector_full = np.array([weights_current.get(t, 0.0) for t in returns.columns], dtype=float)
                if weights_vector_full.sum() <= 0:
                    st.error("A carteira atual não possui pesos válidos após alinhamento com os dados baixados.")
                else:
                    weights_vector = weights_vector_full / weights_vector_full.sum()
                    portfolio_current = returns @ weights_vector

                    cov_matrix = np.asarray(252 * returns.cov())
                    annualized_returns_vec = np.asarray(252 * returns.mean())
                    rent_alvo = np.arange(0.01, 1.2, 0.0025)

                    def port_ret(w):
                        return w @ annualized_returns_vec

                    def port_rets_series(w):
                        return returns @ np.asarray(w)

                    def port_vol(w):
                        return np.sqrt(w @ cov_matrix @ w)

                    def port_semivariance(w):
                        pr = port_rets_series(w)
                        return 252 * np.mean(np.square(np.minimum(pr, 0)))

                    def port_mad(w):
                        pr = port_rets_series(w)
                        return np.mean(np.abs(pr - np.mean(pr))) * np.sqrt(252)

                    def port_cvar(w, alpha):
                        pr = port_rets_series(w)
                        if len(pr) == 0:
                            return np.nan
                        thr = np.quantile(pr, 1 - alpha)
                        tail = pr[pr <= thr]
                        return -thr if len(tail) == 0 else -tail.mean()

                    def port_risk(w):
                        if risk_measure_comparador == "Mean-Variance":
                            return port_vol(w)
                        elif risk_measure_comparador == "Mean-Semivariance":
                            return np.sqrt(port_semivariance(w))
                        elif risk_measure_comparador == "Mean Absolute Deviation (MAD)":
                            return port_mad(w)
                        elif risk_measure_comparador == "CVaR / Expected Shortfall":
                            return port_cvar(w, confidence_level_comparador)
                        return port_vol(w)

                    g_min_t1    = st.session_state.get("alloc_g_min", 0.0)
                    g_max_t1    = st.session_state.get("alloc_g_max", 1.0)
                    per_asset_t1 = st.session_state.get("alloc_per_asset", {})
                    groups_t1   = st.session_state.get("alloc_groups", [])

                    n_assets  = len(returns.columns)
                    bounds_t1 = [per_asset_t1.get(t, (g_min_t1, g_max_t1)) for t in returns.columns]

                    def minimize_risk_for_target(r):
                        base_cons = [
                            {"type": "eq", "fun": lambda w: np.sum(w) - 1},
                            {"type": "eq", "fun": lambda w: port_ret(w) - r},
                        ]
                        grp_cons = _build_opt_constraints(list(returns.columns), groups_t1)
                        res = minimize(port_risk, x0=[1 / n_assets] * n_assets, bounds=bounds_t1, constraints=base_cons + grp_cons)
                        if res.success:
                            return port_risk(res.x), r, res.x

                    frontier_points = [minimize_risk_for_target(r) for r in rent_alvo]
                    frontier_points = [pt for pt in frontier_points if pt is not None]
                    frontier_df = pd.DataFrame(frontier_points, columns=["Risk", "Return", "Weights"])
                    prog.progress(80, text="Executando otimização de carteira...")

                    if frontier_df.empty:
                        st.error(
                            "Não foi possível construir a fronteira eficiente. "
                            "Verifique as restrições de alocação — os limites mínimos/máximos "
                            "podem estar incompatíveis com os ativos selecionados."
                        )
                        prog.empty()
                        st.session_state["comparison_running"] = False
                    else:
                        frontier_df["Sharpe"] = frontier_df["Return"] / frontier_df["Risk"]

                        # A curva contém as duas pernas da "bala" de Markowitz: o ramo
                        # eficiente (retorno cresce com o risco) e o ramo dominado
                        # (mesmo risco, retorno menor). Sem filtrar isso, o ponto de
                        # "Mesmo Risco"/"Mesmo Retorno" mais próximo por distância pode
                        # cair no ramo dominado. Extrai apenas o envelope superior
                        # (Pareto-eficiente): ordena por risco crescente e mantém só os
                        # pontos que batem um novo recorde de retorno.
                        _sorted_by_risk = frontier_df.sort_values("Risk").reset_index(drop=True)
                        _cummax_ret = _sorted_by_risk["Return"].cummax()
                        efficient_df = _sorted_by_risk[_sorted_by_risk["Return"] >= _cummax_ret].reset_index(drop=True)

                        target_return = port_ret(weights_vector)
                        target_risk = port_risk(weights_vector)
                        same_risk_opt = efficient_df.iloc[(efficient_df["Risk"] - target_risk).abs().argmin()]
                        same_return_opt = efficient_df.iloc[(efficient_df["Return"] - target_return).abs().argmin()]

                        prog.progress(82, text="Gerando gráfico da fronteira...")
                        st.markdown("### Selecione a Carteira Otimizada para Backtesting")
                        choice = st.radio(
                            "Selecione a Carteira:",
                            ["Mesmo Risco (↑ Retorno)", "Mesmo Retorno (↓ Risco)", "Personalizado (fronteira)"],
                            key="portfolio_choice_comparador",
                        )

                        if choice == "Mesmo Risco (↑ Retorno)":
                            ponto_sel = same_risk_opt
                            label_opt = "Otimizado - Mesmo Risco"
                        elif choice == "Mesmo Retorno (↓ Risco)":
                            ponto_sel = same_return_opt
                            label_opt = "Otimizado - Mesmo Retorno"
                        else:
                            risco_sel = st.slider(
                                "Nível de risco da carteira otimizada:",
                                min_value=float(efficient_df["Risk"].min()),
                                max_value=float(efficient_df["Risk"].max()),
                                value=float(same_risk_opt["Risk"]),
                                format="%.4f",
                                key="risk_slider_comparador",
                            )
                            ponto_sel = efficient_df.iloc[(efficient_df["Risk"] - risco_sel).abs().argmin()]
                            label_opt = "Otimizado - Personalizado"

                        selected_weights = ponto_sel["Weights"]

                        st.caption(f"Método de construção da fronteira: {risk_measure_comparador}")
                        st.markdown("### Fronteira Eficiente & Comparação de Carteiras")
                        fig_f = go.Figure()
                        fig_f.add_trace(go.Scatter(x=efficient_df["Risk"], y=efficient_df["Return"], mode="lines", name="Fronteira Eficiente"))
                        fig_f.add_trace(go.Scatter(x=[target_risk], y=[target_return], mode="markers+text", marker=dict(color="red", size=12), name="Carteira Atual", text=["Atual"], textposition="top center"))
                        fig_f.add_trace(go.Scatter(x=[same_risk_opt["Risk"]], y=[same_risk_opt["Return"]], mode="markers+text", marker=dict(color="green", size=10, symbol="diamond"), name="Mesmo Risco", text=["↑ Retorno"], textposition="bottom center"))
                        fig_f.add_trace(go.Scatter(x=[same_return_opt["Risk"]], y=[same_return_opt["Return"]], mode="markers+text", marker=dict(color="blue", size=10, symbol="diamond"), name="Mesmo Retorno", text=["↓ Risco"], textposition="bottom center"))
                        if choice == "Personalizado (fronteira)":
                            fig_f.add_trace(go.Scatter(x=[ponto_sel["Risk"]], y=[ponto_sel["Return"]], mode="markers+text", marker=dict(color="orange", size=13, symbol="star"), name="Selecionado", text=["★"], textposition="top center"))
                        fig_f.update_layout(
                            title=f"Fronteira Eficiente & Comparação de Carteiras ({risk_measure_comparador})",
                            xaxis_title=f"Risco ({risk_measure_comparador})",
                            yaxis_title="Retorno Esperado",
                        )
                        st.plotly_chart(fig_f, use_container_width=True)

                        # Versão limpa da fronteira para o relatório PPTX
                        _fig_frontier_report = go.Figure()
                        _fig_frontier_report.add_trace(go.Scatter(
                            x=efficient_df["Risk"], y=efficient_df["Return"],
                            mode="lines", name="Fronteira Eficiente",
                            line=dict(color=_REPORT_PALETTE[0], width=2),
                        ))
                        _fig_frontier_report.add_trace(go.Scatter(
                            x=[target_risk], y=[target_return],
                            mode="markers+text", name="Carteira Atual",
                            marker=dict(color="#DC2626", size=12),
                            text=["Atual"], textposition="top center",
                        ))
                        _fig_frontier_report.add_trace(go.Scatter(
                            x=[same_risk_opt["Risk"]], y=[same_risk_opt["Return"]],
                            mode="markers+text", name="Mesmo Risco (↑ Retorno)",
                            marker=dict(color="#059669", size=10, symbol="diamond"),
                            text=["↑ Retorno"], textposition="bottom center",
                        ))
                        _fig_frontier_report.add_trace(go.Scatter(
                            x=[same_return_opt["Risk"]], y=[same_return_opt["Return"]],
                            mode="markers+text", name="Mesmo Retorno (↓ Risco)",
                            marker=dict(color="#1D4ED8", size=10, symbol="diamond"),
                            text=["↓ Risco"], textposition="bottom center",
                        ))
                        _fig_frontier_report.add_trace(go.Scatter(
                            x=[ponto_sel["Risk"]], y=[ponto_sel["Return"]],
                            mode="markers+text", name="Carteira Otimizada",
                            marker=dict(color=_REPORT_PALETTE[1], size=14, symbol="star"),
                            text=["Otimizada"], textposition="top center",
                        ))
                        _fig_frontier_report.update_layout(**_yc_layout(
                            xaxis_title="Risco",
                            yaxis_title="Retorno Esperado",
                        ))
                        st.session_state["t1_frontier_fig"] = _fig_frontier_report

                        st.subheader("Carteira Atual × Sugerida")
                        patrimonio = st.number_input(
                            "Patrimônio total (R$):",
                            min_value=0.0,
                            value=1_000_000.0,
                            step=10_000.0,
                            format="%.2f",
                            key="portfolio_patrimonio_t1",
                        )

                        pesos_df = pd.DataFrame({"Sugerida": selected_weights}, index=returns.columns)
                        mov_rows = []
                        for ticker in returns.columns:
                            w_atual = weights_current.get(ticker, 0.0)
                            w_sug   = float(pesos_df.loc[ticker, "Sugerida"])
                            if w_atual > 0.001 or w_sug > 0.001:
                                mov_rows.append({
                                    "Ativo":           symbol_info.get(ticker, ticker),
                                    "Ticker":          ticker,
                                    "Atual (%)":       w_atual,
                                    "Sugerida (%)":    w_sug,
                                    "Diferença (%)":   w_sug - w_atual,
                                    "Financeiro (R$)": (w_sug - w_atual) * patrimonio,
                                })
                        mov_df = (
                            pd.DataFrame(mov_rows)
                            .sort_values("Sugerida (%)", ascending=False)
                            .reset_index(drop=True)
                        )

                        def _color_diff(v):
                            if isinstance(v, float):
                                if v > 0.0005:
                                    return "color:#059669;font-weight:bold"
                                if v < -0.0005:
                                    return "color:#dc2626"
                            return ""

                        st.dataframe(
                            mov_df.style
                                  .format({
                                      "Atual (%)":       "{:.2%}",
                                      "Sugerida (%)":    "{:.2%}",
                                      "Diferença (%)":   "{:+.2%}",
                                      "Financeiro (R$)": _brl,
                                  })
                                  .map(_color_diff, subset=["Diferença (%)"]),
                            use_container_width=True,
                            hide_index=True,
                        )

                        returns_opt = returns @ selected_weights
                        prog.progress(85, text="Preparando séries de comparação...")

                        common_idx = portfolio_current.index.intersection(returns_opt.index)
                        if not benchmark_returns.empty:
                            common_idx = common_idx.intersection(benchmark_returns.index)
                        portfolio_current = portfolio_current.loc[common_idx]
                        returns_opt = returns_opt.loc[common_idx]
                        if not benchmark_returns.empty:
                            benchmark_returns = pd.Series(benchmark_returns).dropna().loc[common_idx]

                        st.markdown("### Comparação de Métricas de Desempenho")
                        rf = _get_rf_rate(currency_choice, start_date, end_date, portfolio_current.index)

                        _m_atual = _calcular_metricas(portfolio_current, rf, "Carteira Atual")
                        _m_opt   = _calcular_metricas(returns_opt, rf, label_opt)
                        _exibir_metricas([_m_atual, _m_opt])

                        series_dict = {"Carteira Atual": portfolio_current, label_opt: returns_opt}
                        if not benchmark_returns.empty:
                            series_dict[benchmark_choice_comparador] = benchmark_returns

                        prog.progress(88, text="Gerando gráficos...")
                        _graficos_comparacao(series_dict)
                        _matriz_correlacao(returns)

                        _tabela_rentabilidade_mensal(
                            {"Carteira Atual": portfolio_current, label_opt: returns_opt},
                            benchmark_series=benchmark_returns if not benchmark_returns.empty else None,
                            benchmark_label=benchmark_choice_comparador if not benchmark_returns.empty else None,
                        )

                        prog.progress(100, text="Comparação concluída.")
                        st.session_state["comparison_running"] = False
                        prog.empty()

                        st.session_state["t1_report_series"]  = series_dict
                        st.session_state["t1_report_returns"] = returns
                        st.session_state["t1_weights_atual"]  = weights_current
                        st.session_state["t1_weights_sug"]    = dict(zip(returns.columns, selected_weights))
                        st.session_state["t1_benchmark"]      = benchmark_choice_comparador
                        st.session_state["t1_frontier_metrics"] = {
                            "atual_ret":         target_return,
                            "atual_risco":       target_risk,
                            "mesmo_risco_ret":   float(same_risk_opt["Return"]),
                            "mesmo_risco_risco": float(same_risk_opt["Risk"]),
                            "mesmo_ret_ret":     float(same_return_opt["Return"]),
                            "mesmo_ret_risco":   float(same_return_opt["Risk"]),
                        }
                        st.session_state["t1_metrics_data"] = {
                            "ret_atual":    _m_atual["retorno_anual"],
                            "ret_opt":      _m_opt["retorno_anual"],
                            "vol_atual":    _m_atual["volatilidade"],
                            "vol_opt":      _m_opt["volatilidade"],
                            "sharpe_atual": _m_atual["sharpe"],
                            "sharpe_opt":   _m_opt["sharpe"],
                            "dd_atual":     _m_atual["max_drawdown"],
                            "dd_opt":       _m_opt["max_drawdown"],
                        }

                _render_report_ui("t1")

        # =============================================================
        # TAB 2 — Comparação Simples
        # =============================================================
        with tab2:
            st.subheader("Comparação Simples de Carteiras")
            st.markdown("Compare a carteira atual com uma carteira sugerida e/ou benchmark, sem otimização.")

            if "simples_ready" not in st.session_state:
                st.session_state["simples_ready"] = False

            if st.session_state["simples_ready"]:
                col_back, _ = st.columns([1, 5])
                with col_back:
                    if st.button("← Voltar à configuração", key="back_to_setup_simples"):
                        st.session_state["simples_ready"] = False
                        st.rerun()

            if not st.session_state["simples_ready"]:
                _render_synth_expander("t2", ticker_options, ticker_lookup)
                _currency_t2 = st.session_state.get("currency_simples", "Real")
                st.markdown("#### Carteira Atual")
                atual_portfolio, weight_sum_atual = _portfolio_editor(
                    default_symbols, default_weights, ticker_options, ticker_lookup, "t2_atual",
                    currency_choice=_currency_t2,
                )
                valida_atual = abs(weight_sum_atual - 100) <= 0.01
                if not valida_atual:
                    st.error("A alocação da carteira atual deve somar 100%.")

                incluir_sugerida = st.checkbox("Incluir carteira sugerida para comparação", key="incluir_sugerida")
                sugerida_portfolio = None
                valida_sugerida = True

                if incluir_sugerida:
                    st.markdown("#### Carteira Sugerida")
                    sugerida_portfolio, weight_sum_sug = _portfolio_editor(
                        default_symbols, default_weights, ticker_options, ticker_lookup, "t2_sugerida",
                        currency_choice=_currency_t2,
                    )
                    valida_sugerida = abs(weight_sum_sug - 100) <= 0.01
                    if not valida_sugerida:
                        st.error("A alocação da carteira sugerida deve somar 100%.")

                benchmark_choice_s = st.selectbox(
                    "Selecione o Benchmark:",
                    ["Nenhum"] + list(benchmarks.keys()),
                    index=0,
                    key="benchmark_simples",
                )
                currency_choice_s = st.selectbox(
                    "Selecione a Moeda", ["Real", "Dollar", "Euro"], index=0, key="currency_simples"
                )

                tickers_s = list(atual_portfolio["Ticker"].dropna().unique())
                if sugerida_portfolio is not None:
                    tickers_s = list(dict.fromkeys(tickers_s + list(sugerida_portfolio["Ticker"].dropna().unique())))

                today_s = pd.Timestamp.today().date()

                _simples_invalido = not valida_atual or not valida_sugerida
                if len(tickers_s) >= 1:
                    _synth_setup_s = st.session_state.get("synthetic_assets", {})
                    yf_tickers_s = [t for t in tickers_s if t not in symbol_br_indices and t not in _synth_setup_s]
                    br_tickers_s = [t for t in tickers_s if t in symbol_br_indices]
                    price_data_s = pd.DataFrame()
                    for sym in yf_tickers_s:
                        close = safe_download(sym, start="2000-01-01", end=today_s)
                        if not close.empty:
                            price_data_s[sym] = close
                        else:
                            st.warning(f"⚠️ Falha ao baixar dados para: {sym}")

                    # Determina intervalo disponível considerando TODOS os ativos em conjunto
                    synth_tickers_setup_s = [t for t in tickers_s if t in _synth_setup_s]
                    combined_setup_s = price_data_s.copy() if not price_data_s.empty else pd.DataFrame()
                    if br_tickers_s:
                        br_idx_setup_s = _get_br_slice(start="2000-01-01", end=str(today_s))
                        for _bt in br_tickers_s:
                            if _bt in br_idx_setup_s.columns:
                                combined_setup_s[_bt] = br_idx_setup_s[_bt]
                    _datas_ok = False
                    if not combined_setup_s.empty:
                        _common = combined_setup_s.dropna()
                        first_s = _common.index.min() if not _common.empty else combined_setup_s.index.min()
                        last_s = _common.index.max() if not _common.empty else combined_setup_s.index.max()
                        _datas_ok = True
                    elif synth_tickers_setup_s:
                        _ref = _synth_setup_s[synth_tickers_setup_s[0]]["ref_ticker"]
                        _ref_data = safe_download(_ref, start="2000-01-01", end=today_s)
                        if not _ref_data.empty:
                            first_s = _ref_data.index.min()
                            last_s = _ref_data.index.max()
                        else:
                            first_s = pd.Timestamp("2010-01-01")
                            last_s = pd.Timestamp(today_s)
                        _datas_ok = True
                    else:
                        st.error("Nenhum dado disponível para os ativos selecionados.")

                    if _datas_ok:
                        col1, col2 = st.columns(2)
                        with col1:
                            start_date_s = st.date_input("Data Inicial", value=first_s.date(), min_value=first_s.date(), max_value=last_s.date(), key="t2_start")
                        with col2:
                            end_date_s = st.date_input("Data Final", value=last_s.date(), min_value=first_s.date(), max_value=last_s.date(), key="t2_end")

                        if st.button("Executar Comparação", key="run_simples", disabled=_simples_invalido):
                            st.session_state["simples_atual"] = atual_portfolio
                            st.session_state["simples_sugerida"] = sugerida_portfolio if incluir_sugerida else None
                            st.session_state["simples_benchmark"] = benchmark_choice_s
                            st.session_state["simples_currency"] = currency_choice_s
                            st.session_state["simples_start"] = start_date_s
                            st.session_state["simples_end"] = end_date_s
                            st.session_state["simples_ready"] = True
                            # Ver comentário equivalente na aba 1 — garante reidratação
                            # correta a partir do snapshot em ciclos repetidos.
                            for _k in ("t2_atual_editor", "t2_atual_editor_val", "t2_sugerida_editor", "t2_sugerida_editor_val"):
                                st.session_state.pop(_k, None)
                            st.rerun()

            else:
                prog = st.progress(0, text="Iniciando comparação...")
                atual_portfolio = st.session_state["simples_atual"]
                sugerida_portfolio = st.session_state["simples_sugerida"]
                benchmark_choice_s = st.session_state["simples_benchmark"]
                currency_choice_s = st.session_state["simples_currency"]
                start_date_s = st.session_state["simples_start"]
                end_date_s = st.session_state["simples_end"]
                today_s = pd.Timestamp.today().date()

                tickers_atual = list(atual_portfolio["Ticker"].dropna().unique())
                tickers_sug = list(sugerida_portfolio["Ticker"].dropna().unique()) if sugerida_portfolio is not None else []
                all_tickers = list(dict.fromkeys(tickers_atual + tickers_sug))

                weights_atual = atual_portfolio.set_index("Ticker")["Weight"].to_dict()
                weights_sug = sugerida_portfolio.set_index("Ticker")["Weight"].to_dict() if sugerida_portfolio is not None else {}
                s_a = sum(weights_atual.values())
                weights_atual = {k: v / s_a for k, v in weights_atual.items()} if s_a > 0 else weights_atual
                if weights_sug:
                    s_s = sum(weights_sug.values())
                    weights_sug = {k: v / s_s for k, v in weights_sug.items()} if s_s > 0 else weights_sug

                prog.progress(10, text="Baixando dados dos ativos...")
                _synth_s = st.session_state.get("synthetic_assets", {})
                yf_tickers_r = [t for t in all_tickers if t not in symbol_br_indices and t not in _synth_s]
                br_tickers_r = [t for t in all_tickers if t in symbol_br_indices]
                synth_tickers_r = [t for t in all_tickers if t in _synth_s]

                raw = yf.download(yf_tickers_r, start=start_date_s, end=end_date_s, group_by="ticker", progress=False) if yf_tickers_r else pd.DataFrame()
                if not raw.empty and isinstance(raw.columns, pd.MultiIndex):
                    prices_s = raw.xs("Close", axis=1, level=1)
                    prices_s = prices_s[[t for t in yf_tickers_r if t in prices_s.columns]].dropna()
                elif not raw.empty:
                    prices_s = raw[[t for t in yf_tickers_r if t in raw.columns]].dropna()
                else:
                    prices_s = pd.DataFrame()

                prog.progress(35, text="Calculando retornos...")
                returns_s = prices_s.pct_change().dropna() if not prices_s.empty else pd.DataFrame()
                returns_s = _ajustar_cambio(returns_s, currency_choice_s, today_s)

                # Mescla retornos dos índices BR
                if br_tickers_r:
                    br_idx_r = _get_br_slice(start=str(start_date_s), end=str(end_date_s))
                    if returns_s.empty:
                        returns_s = br_idx_r[[t for t in br_tickers_r if t in br_idx_r.columns]].dropna()
                    else:
                        for t in br_tickers_r:
                            if t in br_idx_r.columns:
                                returns_s[t] = br_idx_r[t]
                        avail_br_r = [t for t in br_tickers_r if t in returns_s.columns]
                        if avail_br_r:
                            returns_s = returns_s.dropna(subset=avail_br_r)

                # Adiciona retornos dos ativos sintéticos
                for _st_name in synth_tickers_r:
                    _si = _synth_s[_st_name]
                    _sr = _compute_synthetic_returns(_si["ref_ticker"], start_date_s, end_date_s)
                    if _sr is not None:
                        returns_s[_st_name] = _si["beta"] * _sr
                    else:
                        st.warning(f"⚠️ Não foi possível calcular retornos para o ativo sintético '{_st_name}' (referência: {_si['ref_ticker']}).")

                benchmark_ticker_s = benchmarks.get(benchmark_choice_s)
                benchmark_returns_s = pd.Series(dtype=float)
                if benchmark_ticker_s is not None and benchmark_ticker_s not in BR_BENCH_TICKERS:
                    prog.progress(50, text="Baixando benchmark...")
                    bm_price_s = safe_download(benchmark_ticker_s, start=start_date_s, end=end_date_s)
                    if not bm_price_s.empty:
                        benchmark_returns_s = pd.Series(bm_price_s).pct_change().dropna()
                        benchmark_returns_s = _ajustar_cambio_benchmark(benchmark_returns_s, benchmark_ticker_s, currency_choice_s, today_s)
                    else:
                        st.warning(f"⚠️ Falha ao baixar dados do benchmark: {benchmark_choice_s}")
                elif benchmark_ticker_s in BR_BENCH_TICKERS:
                    prog.progress(50, text="Baixando benchmark...")
                    br_idx_bm_s = _get_br_slice(start=str(start_date_s), end=str(end_date_s))
                    if benchmark_ticker_s in br_idx_bm_s.columns:
                        benchmark_returns_s = br_idx_bm_s[benchmark_ticker_s].dropna()
                    elif benchmark_ticker_s == "CDI":
                        try:
                            benchmark_returns_s = _load_cdi_bcb(start=str(start_date_s), end=str(end_date_s))
                        except Exception as _cdi_e:
                            st.warning(f"⚠️ Não foi possível carregar o CDI do Banco Central: {_cdi_e}")
                    else:
                        st.warning(f"⚠️ Índice '{benchmark_choice_s}' não disponível no período selecionado.")

                prog.progress(65, text="Computando séries de carteiras...")
                w_vec_atual = np.array([weights_atual.get(t, 0.0) for t in returns_s.columns], dtype=float)
                if w_vec_atual.sum() > 0:
                    w_vec_atual /= w_vec_atual.sum()
                ret_atual = returns_s @ w_vec_atual

                series_dict = {"Carteira Atual": ret_atual}

                if weights_sug:
                    w_vec_sug = np.array([weights_sug.get(t, 0.0) for t in returns_s.columns], dtype=float)
                    if w_vec_sug.sum() > 0:
                        w_vec_sug /= w_vec_sug.sum()
                        series_dict["Carteira Sugerida"] = returns_s @ w_vec_sug

                if not benchmark_returns_s.empty:
                    series_dict[benchmark_choice_s] = benchmark_returns_s

                prog.progress(75, text="Calculando métricas...")
                rf_s = _get_rf_rate(currency_choice_s, start_date_s, end_date_s, ret_atual.index)

                st.markdown("### Comparação de Métricas de Desempenho")
                _ms_list = [
                    _calcular_metricas(s, rf_s, lbl)
                    for lbl, s in series_dict.items()
                    if lbl != benchmark_choice_s
                ]
                _exibir_metricas(_ms_list)

                prog.progress(88, text="Gerando gráficos...")
                _graficos_comparacao(series_dict)
                _matriz_correlacao(returns_s)

                portfolios_tabela = {k: v for k, v in series_dict.items() if k != benchmark_choice_s}
                _tabela_rentabilidade_mensal(
                    portfolios_tabela,
                    benchmark_series=benchmark_returns_s if not benchmark_returns_s.empty else None,
                    benchmark_label=benchmark_choice_s if not benchmark_returns_s.empty else None,
                )

                prog.progress(100, text="Comparação concluída.")
                prog.empty()

                st.session_state["t2_report_series"]  = series_dict
                st.session_state["t2_report_returns"] = returns_s
                st.session_state["t2_weights_atual"]  = weights_atual
                st.session_state["t2_weights_sug"]    = weights_sug if weights_sug else None
                st.session_state["t2_benchmark"]      = benchmark_choice_s
                if len(_ms_list) >= 2:
                    st.session_state["t2_metrics_data"] = {
                        "ret_atual":    _ms_list[0]["retorno_anual"],
                        "ret_opt":      _ms_list[1]["retorno_anual"],
                        "vol_atual":    _ms_list[0]["volatilidade"],
                        "vol_opt":      _ms_list[1]["volatilidade"],
                        "sharpe_atual": _ms_list[0]["sharpe"],
                        "sharpe_opt":   _ms_list[1]["sharpe"],
                        "dd_atual":     _ms_list[0]["max_drawdown"],
                        "dd_opt":       _ms_list[1]["max_drawdown"],
                    }

            _render_report_ui("t2")

    except Exception as e:
        st.error("Ocorreu um erro inesperado durante a comparação de carteiras.")
        st.text(f"Detalhes técnicos: {e}")


# Backward compatibility, in case main.py still calls comparador()
comparador = show_comparador
